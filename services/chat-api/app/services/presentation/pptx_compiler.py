"""PPTX compiler — converts FreeformDeckBlueprint to a .pptx file.

Consumes the same upstream blueprint as html_renderer, using shared
render_utils for coordinate conversion, z-ordering, style resolution
and chart spec extraction.
"""
from __future__ import annotations

import logging
import re
import tempfile
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from app.services.presentation.contracts import (
    FreeformBlock,
    FreeformDeckBlueprint,
    FreeformPageBlueprint,
)
from app.services.presentation.render_utils import (
    SLIDE_H_EMU,
    SLIDE_W_EMU,
    chart_theme_from_deck,
    infer_text_vertical_align,
    norm_to_emu,
    parse_linear_gradient,
    render_order,
    resolve_block_content,
    resolve_chart_spec,
    resolve_color_rgb,
    resolve_font_size_pt,
    resolve_font_weight,
    resolve_line_endpoints,
    resolve_page_coords,
    resolve_text_align,
    resolve_vertical_align,
)
from app.services.presentation.style_contract import PresentationStyleContract
from app.services.presentation.image_layout import crop_image_bytes_to_aspect
from app.services.presentation.text_box_style import parse_css_padding
from app.utils.oss_uploader import AliyunOSSUploader

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _rgb(value: str) -> RGBColor:
    """Best-effort parse any CSS colour string to RGBColor."""
    rgb = resolve_color_rgb(value)
    if rgb is None:
        return RGBColor(0x0F, 0x17, 0x2A)
    return RGBColor(rgb[0], rgb[1], rgb[2])


def _apply_fill_alpha(fill, alpha: float) -> None:
    """Inject <a:alpha val=".."/> into a solid or gradient fill's color nodes.

    python-pptx exposes no transparency API, so we edit OOXML directly.
    alpha is 0..1 where 1=opaque; OOXML alpha val is in 1000ths of a percent.
    """
    val = max(0, min(100000, int(round(alpha * 100000))))
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    fill_elem = fill._xPr.find(qn("a:solidFill"))
    color_nodes = []
    if fill_elem is not None:
        for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr", "a:sysClr"):
            node = fill_elem.find(qn(tag))
            if node is not None:
                color_nodes.append(node)
    else:
        grad = fill._xPr.find(qn("a:gradFill"))
        if grad is not None:
            for gs in grad.findall(qn("a:gsLst") + "/" + qn("a:gs")):
                for tag in ("a:srgbClr", "a:schemeClr", "a:prstClr", "a:sysClr"):
                    node = gs.find(qn(tag))
                    if node is not None:
                        color_nodes.append(node)
    for node in color_nodes:
        existing = node.find(qn("a:alpha"))
        if existing is not None:
            node.remove(existing)
        alpha_elem = etree.SubElement(node, f"{{{a_ns}}}alpha")
        alpha_elem.set("val", str(val))


_PP_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# ---------------------------------------------------------------------------
# Compiler
# ---------------------------------------------------------------------------


class PptxCompiler:
    """Render a FreeformDeckBlueprint into a .pptx file."""

    def compile(self, blueprint: FreeformDeckBlueprint) -> bytes:
        """Return raw .pptx bytes."""
        blueprint = PresentationStyleContract().canonicalize(blueprint)
        prs = Presentation()
        prs.slide_width = Emu(SLIDE_W_EMU)
        prs.slide_height = Emu(SLIDE_H_EMU)
        blank_layout = prs.slide_layouts[6]

        for page in list(blueprint.pages or []):
            slide = prs.slides.add_slide(blank_layout)
            self._render_page(slide, page, blueprint)

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    async def compile_and_upload(
        self,
        *,
        blueprint: FreeformDeckBlueprint,
        user_id: str,
        filename: str = "",
    ) -> Dict[str, Any]:
        """Compile blueprint → .pptx → upload to OSS → return metadata."""
        pptx_bytes = self.compile(blueprint)
        safe_name = filename.strip() or f"{self._safe_filename(blueprint.deck_goal or blueprint.deck_id)}.pptx"
        if not safe_name.endswith(".pptx"):
            safe_name += ".pptx"

        uploader = AliyunOSSUploader()
        _url, object_path = uploader.upload_bytes_with_path(
            pptx_bytes,
            user_id=user_id,
            file_name=safe_name,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        signed_url = uploader.sign_url(object_path)
        logger.info(
            "pptx_v5_uploaded deck_id=%s slides=%s size_kb=%s",
            blueprint.deck_id,
            len(blueprint.pages),
            len(pptx_bytes) // 1024,
        )
        return {
            "type": "pptx",
            "url": signed_url,
            "filename": safe_name,
            "title": blueprint.deck_goal or "Presentation",
            "object_path": object_path,
            "slide_count": len(blueprint.pages),
        }

    # ------------------------------------------------------------------
    # Page rendering
    # ------------------------------------------------------------------

    def _render_page(
        self,
        slide,
        page: FreeformPageBlueprint,
        deck: FreeformDeckBlueprint,
    ) -> None:
        # Page background
        self._apply_page_background(slide, page, deck)

        # Sort blocks by z-order then render
        sorted_blocks = sorted(list(page.blocks or []), key=render_order)
        for block in sorted_blocks:
            self._render_block(slide, block, deck)

    def _apply_page_background(
        self, slide, page: FreeformPageBlueprint, deck: FreeformDeckBlueprint,
    ) -> None:
        page_style = dict(page.style or {})
        bg_value = str(
            page_style.get("background")
            or page_style.get("background_color")
            or deck.theme.page_background
            or "#ffffff"
        ).strip()

        gradient = parse_linear_gradient(bg_value)
        if gradient:
            self._apply_gradient_fill(slide.background.fill, gradient)
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(bg_value)

    # ------------------------------------------------------------------
    # Block dispatch
    # ------------------------------------------------------------------

    def _render_block(
        self,
        slide,
        block: FreeformBlock,
        deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float] = (0.0, 0.0, 1.0, 1.0),
    ) -> None:
        block_type = str(block.type or "text_box").strip().lower()
        dispatch = {
            "rectangle": self._render_rectangle,
            "circle": self._render_circle,
            "text_box": self._render_text_box,
            "line": self._render_line,
            "group": self._render_group,
            "icon": self._render_icon,
            "chart": self._render_chart,
            "image": self._render_image,
        }
        handler = dispatch.get(block_type, self._render_text_box)
        handler(slide, block, deck, parent_rect)

    # ------------------------------------------------------------------
    # Shape renderers
    # ------------------------------------------------------------------

    def _render_rectangle(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        style = dict(block.style or {})

        border_radius = self._parse_border_radius(style)
        if border_radius > 0:
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height,
            )
            self._set_corner_radius(shape, border_radius, width, height)
        else:
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height,
            )

        self._apply_shape_fill(shape, style, deck)
        self._apply_shape_border(shape, style)

        # Render children inside the rectangle (e.g. text inside a card)
        current_rect = (x, y, w, h)
        sorted_children = sorted(list(block.children or []), key=render_order)
        for child in sorted_children:
            self._render_block(slide, child, deck, current_rect)

    def _render_circle(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        style = dict(block.style or {})

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height,
        )
        self._apply_shape_fill(shape, style, deck)
        self._apply_shape_border(shape, style)
        content = resolve_block_content(block)
        if content:
            circle_style = dict(style)
            circle_style.setdefault("text_align", "center")
            circle_style.setdefault("vertical_align", "middle")
            self._write_text_frame(shape.text_frame, content, circle_style, deck, box_height_norm=h)

        current_rect = (x, y, w, h)
        sorted_children = sorted(list(block.children or []), key=render_order)
        for child in sorted_children:
            self._render_block(slide, child, deck, current_rect)

    def _render_text_box(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        style = dict(block.style or {})
        content = resolve_block_content(block)

        # If the block has a background, render as a shape with text frame
        bg = style.get("background") or style.get("background_color")
        border_radius = self._parse_border_radius(style)
        if bg and str(bg).strip().lower() not in {"transparent", "none", "inherit"}:
            if border_radius > 0:
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height,
                )
                self._set_corner_radius(shape, border_radius, width, height)
            else:
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height,
                )
            self._apply_shape_fill(shape, style, deck)
            self._apply_shape_border(shape, style)
            tf = shape.text_frame
        else:
            box = slide.shapes.add_textbox(left, top, width, height)
            box.line.fill.background()
            tf = box.text_frame

        self._draw_side_borders(slide, style, left, top, width, height)

        self._write_text_frame(tf, content, style, deck, box_height_norm=h)

        current_rect = (x, y, w, h)
        sorted_children = sorted(list(block.children or []), key=render_order)
        for child in sorted_children:
            self._render_block(slide, child, deck, current_rect)

    def _render_line(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x1, y1, x2, y2 = resolve_line_endpoints(block, parent_rect)
        style = dict(block.style or {})

        left_emu = int(min(x1, x2) * SLIDE_W_EMU)
        top_emu = int(min(y1, y2) * SLIDE_H_EMU)
        right_emu = int(max(x1, x2) * SLIDE_W_EMU)
        bottom_emu = int(max(y1, y2) * SLIDE_H_EMU)
        width_emu = max(1, right_emu - left_emu)
        height_emu = max(1, bottom_emu - top_emu)

        # Use a thin rectangle to simulate a line — more reliable than connector
        color_str = str(
            style.get("color")
            or style.get("background")
            or style.get("border_color")
            or getattr(deck.theme, "accent_color", "")
            or "#000000"
        ).strip()

        line_weight = style.get("line_weight")
        try:
            thickness = max(1.0, float(line_weight)) if line_weight not in (None, "") else 2.0
        except Exception:
            thickness = 2.0

        # For diagonal lines, use a freeform; for axis-aligned, use a thin rect
        if abs(x2 - x1) < 0.001 or abs(y2 - y1) < 0.001:
            # Axis-aligned line
            if abs(x2 - x1) < 0.001:
                # Vertical line
                line_w = max(1, int(Pt(thickness)))
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    int(x1 * SLIDE_W_EMU) - line_w // 2, top_emu,
                    line_w, height_emu,
                )
            else:
                # Horizontal line
                line_h = max(1, int(Pt(thickness)))
                shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    left_emu, int(y1 * SLIDE_H_EMU) - line_h // 2,
                    width_emu, line_h,
                )
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(color_str)
            shape.line.fill.background()
        else:
            # Diagonal: use connector
            connector = slide.shapes.add_connector(
                1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                int(x1 * SLIDE_W_EMU), int(y1 * SLIDE_H_EMU),
                int(x2 * SLIDE_W_EMU), int(y2 * SLIDE_H_EMU),
            )
            connector.line.color.rgb = _rgb(color_str)
            connector.line.width = Pt(thickness)

    def _render_group(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x, y, w, h = resolve_page_coords(block, parent_rect)
        style = dict(block.style or {})

        # If the group itself has a visible background, render a backing shape
        bg = style.get("background") or style.get("background_color")
        if bg and str(bg).strip().lower() not in {"transparent", "none", "inherit"}:
            left, top, width, height = norm_to_emu(x, y, w, h)
            border_radius = self._parse_border_radius(style)
            if border_radius > 0:
                backing = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height,
                )
                self._set_corner_radius(backing, border_radius, width, height)
            else:
                backing = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height,
                )
            self._apply_shape_fill(backing, style, deck)
            self._apply_shape_border(backing, style)
            self._draw_side_borders(slide, style, left, top, width, height)
        else:
            # Even without a background, side borders may still be set
            left, top, width, height = norm_to_emu(x, y, w, h)
            self._draw_side_borders(slide, style, left, top, width, height)

        # Recurse into children
        current_rect = (x, y, w, h)
        sorted_children = sorted(list(block.children or []), key=render_order)
        for child in sorted_children:
            self._render_block(slide, child, deck, current_rect)

    def _render_icon(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        from app.services.presentation.svg_to_shape import render_icon_to_slide

        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        style = dict(block.style or {})
        color_str = str(
            style.get("color")
            or getattr(deck.theme, "accent_color", "")
            or "#000000"
        ).strip()

        # Icons are square (24×24 viewBox) — fit uniformly, centered in block area
        icon_size = min(width, height)
        icon_left = left + (width - icon_size) // 2
        icon_top = top + (height - icon_size) // 2

        render_icon_to_slide(
            slide=slide,
            block=block,
            left=icon_left, top=icon_top, width=icon_size, height=icon_size,
            color_rgb=_rgb(color_str),
        )

    def _render_chart(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        from app.services.presentation.pptx_chart_renderer import render_chart_to_slide

        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        chart_type, chart_payload = resolve_chart_spec(block)
        theme = chart_theme_from_deck(deck)

        if chart_type:
            render_chart_to_slide(
                slide=slide,
                chart_type=chart_type,
                payload=chart_payload,
                theme=theme,
                left=left, top=top, width=width, height=height,
            )
        else:
            # Fallback placeholder
            fallback_text = str(chart_payload.get("title") or "Chart")
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height,
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
            shape.line.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
            self._write_text_frame(
                shape.text_frame, fallback_text,
                {"color": "#64748b", "font_size": 14, "text_align": "center"},
                deck,
            )

    def _render_image(
        self, slide, block: FreeformBlock, deck: FreeformDeckBlueprint,
        parent_rect: Tuple[float, float, float, float],
    ) -> None:
        x, y, w, h = resolve_page_coords(block, parent_rect)
        left, top, width, height = norm_to_emu(x, y, w, h)
        style = dict(block.style or {})
        content = str(block.content or "").strip()

        # Try to download and embed a real image URL
        if content and (content.startswith("http://") or content.startswith("https://")):
            image_bytes = self._download_image(content)
            if image_bytes:
                if str(style.get("fit") or style.get("object_fit") or "cover").strip().lower() == "cover":
                    image_bytes = crop_image_bytes_to_aspect(
                        image_bytes,
                        target_width=max(1, int(width)),
                        target_height=max(1, int(height)),
                    )
                slide.shapes.add_picture(BytesIO(image_bytes), left, top, width, height)
                return

        # Placeholder
        prompt_text = str(block.image_prompt or content or "Image").strip()
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF0)
        shape.line.color.rgb = RGBColor(0xB0, 0xB8, 0xC4)
        self._write_text_frame(
            shape.text_frame, prompt_text,
            {"color": "#8896a7", "font_size": 13, "text_align": "center", "vertical_align": "middle"},
            deck,
        )

    # ------------------------------------------------------------------
    # Fill / border / gradient helpers
    # ------------------------------------------------------------------

    def _apply_shape_fill(
        self, shape, style: Dict[str, Any], deck: FreeformDeckBlueprint,
    ) -> None:
        bg = str(
            style.get("background")
            or style.get("background_color")
            or ""
        ).strip()
        if not bg or bg.lower() in {"transparent", "none", "inherit"}:
            shape.fill.background()
            return

        gradient = parse_linear_gradient(bg)
        if gradient:
            self._apply_gradient_fill(shape.fill, gradient)
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(bg)

        # Opacity — python-pptx has no FillFormat.transparency, so inject
        # <a:alpha val="..."/> directly into the fill element's color node.
        # Source of alpha: explicit style.opacity, OR the alpha channel of an
        # rgba(...) background (which resolve_color_rgb silently discards).
        alpha: float | None = None
        opacity = style.get("opacity")
        if opacity is not None:
            try:
                alpha = float(opacity)
            except (TypeError, ValueError):
                alpha = None
        if alpha is None:
            m = re.match(
                r"\s*rgba\(\s*\d+\s*,\s*\d+\s*,\s*\d+\s*,\s*([0-9.]+)\s*\)",
                bg,
            )
            if m:
                try:
                    alpha = float(m.group(1))
                except ValueError:
                    alpha = None
        if alpha is not None and 0.0 <= alpha < 1.0:
            _apply_fill_alpha(shape.fill, alpha)

    def _apply_gradient_fill(self, fill, gradient: Tuple[float, list]) -> None:
        angle, stops = gradient
        fill.gradient()
        fill.gradient_angle = angle

        # python-pptx pre-creates exactly 2 gradient stops and exposes no public
        # API to add more. If we have >2 stops we manipulate the underlying XML
        # (<a:gsLst>) directly; otherwise use the high-level API.
        gs = fill.gradient_stops
        stop_count = len(stops)
        if stop_count <= len(gs):
            for idx, (position, rgb) in enumerate(stops):
                gs[idx].position = position
                gs[idx].color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            return

        # Multi-stop gradient: rebuild <a:gsLst> children from scratch via lxml.
        gs_lst = getattr(gs, "_gsLst", None)
        if gs_lst is None:
            # fallback: apply only the first two stops
            for idx, (position, rgb) in enumerate(stops[:len(gs)]):
                gs[idx].position = position
                gs[idx].color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
            return
        # Clear existing <a:gs> entries and re-emit with correct positions/colors.
        for child in list(gs_lst):
            gs_lst.remove(child)
        for position, rgb in stops:
            pos_int = max(0, min(100000, int(round(float(position) * 100000))))
            gs_el = etree.SubElement(gs_lst, qn("a:gs"), pos=str(pos_int))
            srgb = etree.SubElement(gs_el, qn("a:srgbClr"))
            srgb.set("val", "{:02X}{:02X}{:02X}".format(int(rgb[0]), int(rgb[1]), int(rgb[2])))

    def _draw_side_borders(
        self, slide, style: Dict[str, Any],
        left: int, top: int, width: int, height: int,
    ) -> None:
        """Render per-side borders (border_left/right/top/bottom) as thin
        filled rectangles on top of the host shape. PowerPoint shapes can
        only carry a uniform border, so side-accent strips are emulated
        exactly like the editor does.
        """
        sides = {
            "border_left":   lambda bw: (left, top, bw, height),
            "border_right":  lambda bw: (left + width - bw, top, bw, height),
            "border_top":    lambda bw: (left, top, width, bw),
            "border_bottom": lambda bw: (left, top + height - bw, width, bw),
        }
        pat = re.compile(
            r"(\d+(?:\.\d+)?)\s*px\s+solid\s+(#[0-9a-fA-F]{3,8}|rgba?\([^)]+\))"
        )
        radius_px = self._parse_border_radius(style)
        for key, geom in sides.items():
            val = style.get(key)
            if not val:
                continue
            m = pat.match(str(val).strip())
            if not m:
                continue
            bw_px = float(m.group(1))
            bw_emu = max(1, int(bw_px * SLIDE_W_EMU / 960))
            lx, ly, lw, lh = geom(bw_emu)
            # Use rounded rect only if the panel has a radius AND the strip is
            # wide/tall enough for the radius to render meaningfully without
            # degenerating into a pill. Cap radius at half the strip's SHORT
            # dimension so it never becomes larger than what fits — matches the
            # editor's fabric rx/ry cap behaviour.
            short_side_emu = min(lw, lh)
            radius_emu = int(radius_px * 12700) if radius_px > 0 else 0
            effective_radius_emu = min(radius_emu, short_side_emu // 2)
            if effective_radius_emu > 0:
                strip = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, lx, ly, lw, lh,
                )
                # ratio is relative to the short side
                try:
                    strip.adjustments[0] = effective_radius_emu / short_side_emu
                except Exception:
                    pass
            else:
                strip = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE, lx, ly, lw, lh,
                )
            strip.fill.solid()
            strip.fill.fore_color.rgb = _rgb(m.group(2))
            strip.line.fill.background()

    def _apply_shape_border(self, shape, style: Dict[str, Any]) -> None:
        border_color = str(
            style.get("border_color")
            or style.get("border")
            or ""
        ).strip()

        # Parse "Npx solid #color" shorthand
        border_match = re.match(
            r"(\d+(?:\.\d+)?)\s*px\s+solid\s+(#[0-9a-fA-F]{3,8}|rgba?\([^)]+\))",
            border_color,
        )
        if border_match:
            shape.line.width = Pt(float(border_match.group(1)))
            shape.line.color.rgb = _rgb(border_match.group(2))
            return

        if border_color and border_color.lower() not in {"none", "transparent"}:
            rgb = resolve_color_rgb(border_color)
            if rgb:
                shape.line.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                border_width = style.get("border_width")
                if border_width is not None:
                    try:
                        shape.line.width = Pt(max(0.5, float(border_width)))
                    except (TypeError, ValueError):
                        pass
                return

        # No border specified — make it invisible
        shape.line.fill.background()

    def _parse_border_radius(self, style: Dict[str, Any]) -> float:
        raw = style.get("border_radius")
        if raw is None:
            return 0.0
        try:
            return max(0.0, float(str(raw).replace("px", "").strip()))
        except (TypeError, ValueError):
            return 0.0

    def _set_corner_radius(self, shape, radius_px: float, width_emu: int, height_emu: int) -> None:
        """Set rounded-rectangle corner size via adjustment value.

        python-pptx exposes adjustments as a ratio of the shorter side.
        """
        short_side = min(width_emu, height_emu)
        if short_side <= 0:
            return
        radius_emu = int(radius_px * 12700)  # 1px ≈ 12700 EMU at 96dpi
        ratio = min(0.5, radius_emu / short_side)
        try:
            shape.adjustments[0] = ratio
        except Exception:
            pass

    # ------------------------------------------------------------------
    # Text rendering
    # ------------------------------------------------------------------

    def _write_text_frame(
        self,
        tf,
        content: str,
        style: Dict[str, Any],
        deck: FreeformDeckBlueprint,
        box_height_norm: float = 0.0,
    ) -> None:
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE

        # Vertical alignment
        v_align = infer_text_vertical_align(
            content=content,
            style=style,
            box_height_norm=box_height_norm,
        )
        from pptx.enum.text import MSO_VERTICAL_ANCHOR
        anchor_map = {
            "top": MSO_VERTICAL_ANCHOR.TOP,
            "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
            "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
        }
        tf.paragraphs[0].alignment = _PP_ALIGN_MAP.get(
            resolve_text_align(style.get("text_align")), PP_ALIGN.LEFT,
        )

        # Padding
        pad_top, pad_right, pad_bottom, pad_left = parse_css_padding(style.get("padding"))
        tf.margin_left = int(pad_left * 12700)
        tf.margin_right = int(pad_right * 12700)
        tf.margin_top = int(pad_top * 12700)
        tf.margin_bottom = int(pad_bottom * 12700)

        # Text styling
        font_size = resolve_font_size_pt(style.get("font_size"))
        bold = resolve_font_weight(style.get("font_weight"))
        color_str = str(style.get("color") or getattr(deck.theme, "body_color", "#222222") or "#222222")
        text_align = resolve_text_align(style.get("text_align"))
        alignment = _PP_ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)
        line_height = style.get("line_height")
        try:
            line_spacing = max(0.8, min(3.0, float(line_height))) if line_height is not None else None
        except (TypeError, ValueError):
            line_spacing = None

        font_family = str(
            style.get("font_family")
            or getattr(deck.theme, "font_family", "")
            or ""
        ).strip()
        primary_font = ""
        if font_family:
            primary_font = font_family.split(",")[0].strip().strip("'\"")

        lines = str(content or "").replace("\\n", "\n").splitlines() or [""]
        for idx, line in enumerate(lines):
            if idx == 0:
                paragraph = tf.paragraphs[0]
            else:
                paragraph = tf.add_paragraph()
            paragraph.alignment = alignment
            if line_spacing is not None:
                paragraph.line_spacing = line_spacing
            run = paragraph.add_run()
            run.text = str(line)
            if font_size:
                run.font.size = Pt(max(6, font_size))
            run.font.bold = bold
            run.font.italic = str(style.get("font_style") or "").strip().lower() == "italic"
            run.font.color.rgb = _rgb(color_str)
            if primary_font:
                run.font.name = primary_font

        # Apply vertical anchor after paragraphs are set
        tf.vertical_anchor = anchor_map.get(v_align, MSO_VERTICAL_ANCHOR.TOP)

    # ------------------------------------------------------------------
    # Image download
    # ------------------------------------------------------------------

    @staticmethod
    def _download_image(url: str, timeout: float = 10.0) -> bytes | None:
        try:
            import httpx
            with httpx.Client(timeout=timeout, follow_redirects=True) as client:
                resp = client.get(url)
                if resp.status_code == 200 and len(resp.content) > 100:
                    return resp.content
        except Exception:
            pass
        return None

    # ------------------------------------------------------------------
    # Filename
    # ------------------------------------------------------------------

    @staticmethod
    def _safe_filename(value: str) -> str:
        token = "".join(
            ch if ch.isalnum() or ch in {"-", "_"} else "_"
            for ch in str(value or "").strip()
        )
        return token[:80] or "presentation"
