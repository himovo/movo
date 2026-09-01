"""Convert Tabler SVG icons to OOXML custGeom shapes for PPTX embedding.

Strategy (方案C): parse SVG ``<path d="...">`` into OOXML ``a:custGeom``
path instructions.  PowerPoint renders these as native vector shapes —
scalable, no rasterisation, no external library dependency.

Tabler icons are all 24×24 viewBox, stroke-based, using a small SVG path
instruction set: M, m, L, l, H, h, V, v, C, c, A, a, Z/z.  We convert
each into the corresponding OOXML drawingML path commands.

OOXML custGeom uses EMU units within a local coordinate space defined by
``a:avLst`` and ``a:pathLst``.  We map the SVG 24×24 space to a 240000
EMU square (10000 EMU per SVG unit), which gives plenty of precision.
"""
from __future__ import annotations

import math
import re
import logging
from typing import List, Tuple
from xml.etree.ElementTree import SubElement

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

from app.services.presentation.contracts import FreeformBlock
from app.services.presentation.icon_library import load_inline_svg_map, resolve_icon_name

logger = logging.getLogger(__name__)

# EMU per SVG unit (24 SVG units → 240000 EMU)
_EMU_PER_SVG = 10000
_SVG_SIZE = 24
_GEOM_SIZE = _SVG_SIZE * _EMU_PER_SVG  # 240000


# ---------------------------------------------------------------------------
# SVG path tokeniser
# ---------------------------------------------------------------------------


def _tokenise_path(d: str) -> List[str]:
    """Split an SVG path ``d`` attribute into command + number tokens."""
    return re.findall(r"[MmLlHhVvCcSsQqTtAaZz]|[-+]?(?:\d+\.?\d*|\.\d+)(?:[eE][-+]?\d+)?", d)


def _pop_float(tokens: List[str], idx: int) -> Tuple[float, int]:
    if idx >= len(tokens):
        return 0.0, idx
    try:
        return float(tokens[idx]), idx + 1
    except ValueError:
        return 0.0, idx


# ---------------------------------------------------------------------------
# SVG path → OOXML path elements
# ---------------------------------------------------------------------------


def _svg_to_emu(val: float) -> int:
    return int(round(val * _EMU_PER_SVG))


def _arc_to_beziers(
    cx: float, cy: float, rx: float, ry: float,
    start_angle: float, sweep_angle: float,
    cos_phi: float, sin_phi: float,
) -> List[Tuple[float, float, float, float, float, float]]:
    """Approximate an elliptical arc with cubic Bezier segments.

    Returns list of (c1x, c1y, c2x, c2y, ex, ey) in original coords.
    """
    segments: List[Tuple[float, float, float, float, float, float]] = []
    n = max(1, int(math.ceil(abs(sweep_angle) / (math.pi / 2))))
    delta = sweep_angle / n

    for i in range(n):
        theta1 = start_angle + i * delta
        theta2 = theta1 + delta
        alpha = 4.0 * math.tan(delta / 4.0) / 3.0

        # Points on unit circle
        cos1, sin1 = math.cos(theta1), math.sin(theta1)
        cos2, sin2 = math.cos(theta2), math.sin(theta2)

        # Control points in ellipse space
        ep1x = rx * cos1
        ep1y = ry * sin1
        ep2x = rx * cos2
        ep2y = ry * sin2

        c1x = ep1x - alpha * rx * sin1
        c1y = ep1y + alpha * ry * cos1
        c2x = ep2x + alpha * rx * sin2
        c2y = ep2y - alpha * ry * cos2

        # Rotate back
        def rotate(px, py):
            return px * cos_phi - py * sin_phi + cx, px * sin_phi + py * cos_phi + cy

        rc1 = rotate(c1x, c1y)
        rc2 = rotate(c2x, c2y)
        re2 = rotate(ep2x, ep2y)
        segments.append((rc1[0], rc1[1], rc2[0], rc2[1], re2[0], re2[1]))

    return segments


def _endpoint_to_center(
    x1: float, y1: float, x2: float, y2: float,
    fa: int, fs: int, rx: float, ry: float, phi: float,
) -> Tuple[float, float, float, float, float, float]:
    """Convert SVG arc endpoint parameterisation to center parameterisation."""
    cos_phi = math.cos(phi)
    sin_phi = math.sin(phi)

    dx = (x1 - x2) / 2.0
    dy = (y1 - y2) / 2.0
    x1p = cos_phi * dx + sin_phi * dy
    y1p = -sin_phi * dx + cos_phi * dy

    # Ensure radii are large enough
    lam = (x1p * x1p) / (rx * rx) + (y1p * y1p) / (ry * ry)
    if lam > 1.0:
        s = math.sqrt(lam)
        rx *= s
        ry *= s

    num = max(0.0, rx * rx * ry * ry - rx * rx * y1p * y1p - ry * ry * x1p * x1p)
    den = rx * rx * y1p * y1p + ry * ry * x1p * x1p
    sq = math.sqrt(num / den) if den > 0 else 0.0
    if fa == fs:
        sq = -sq

    cxp = sq * rx * y1p / ry
    cyp = -sq * ry * x1p / rx

    cx = cos_phi * cxp - sin_phi * cyp + (x1 + x2) / 2.0
    cy = sin_phi * cxp + cos_phi * cyp + (y1 + y2) / 2.0

    def angle(ux, uy, vx, vy):
        dot = ux * vx + uy * vy
        mag = math.sqrt(ux * ux + uy * uy) * math.sqrt(vx * vx + vy * vy)
        cos_a = max(-1.0, min(1.0, dot / mag)) if mag > 0 else 1.0
        a = math.acos(cos_a)
        if ux * vy - uy * vx < 0:
            a = -a
        return a

    theta1 = angle(1, 0, (x1p - cxp) / rx, (y1p - cyp) / ry)
    dtheta = angle(
        (x1p - cxp) / rx, (y1p - cyp) / ry,
        (-x1p - cxp) / rx, (-y1p - cyp) / ry,
    )
    if fs == 0 and dtheta > 0:
        dtheta -= 2 * math.pi
    elif fs == 1 and dtheta < 0:
        dtheta += 2 * math.pi

    return cx, cy, rx, ry, theta1, dtheta


def svg_paths_to_ooxml(
    path_strings: List[str],
    stroke_width: float = 2.0,
) -> etree._Element:
    """Convert list of SVG path ``d`` strings into an ``a:custGeom`` element.

    The returned element can replace a shape's ``a:prstGeom`` in the OOXML
    tree to produce a custom vector shape.
    """
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    cust_geom = etree.Element(qn("a:custGeom"))
    etree.SubElement(cust_geom, qn("a:avLst"))
    gdLst = etree.SubElement(cust_geom, qn("a:gdLst"))
    ahLst = etree.SubElement(cust_geom, qn("a:ahLst"))
    cxnLst = etree.SubElement(cust_geom, qn("a:cxnLst"))
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "0")
    rect.set("t", "0")
    rect.set("r", str(_GEOM_SIZE))
    rect.set("b", str(_GEOM_SIZE))

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))

    for d_str in path_strings:
        path_el = etree.SubElement(path_lst, qn("a:path"))
        path_el.set("w", str(_GEOM_SIZE))
        path_el.set("h", str(_GEOM_SIZE))
        path_el.set("stroke", "1")
        path_el.set("fill", "none")

        tokens = _tokenise_path(d_str)
        i = 0
        cur_x, cur_y = 0.0, 0.0
        start_x, start_y = 0.0, 0.0

        while i < len(tokens):
            cmd = tokens[i]
            i += 1

            if cmd in ("M", "m"):
                x, i = _pop_float(tokens, i)
                y, i = _pop_float(tokens, i)
                if cmd == "m":
                    x += cur_x
                    y += cur_y
                cur_x, cur_y = x, y
                start_x, start_y = x, y
                move = etree.SubElement(path_el, qn("a:moveTo"))
                pt = etree.SubElement(move, qn("a:pt"))
                pt.set("x", str(_svg_to_emu(x)))
                pt.set("y", str(_svg_to_emu(y)))
                # Implicit lineTo for subsequent coordinate pairs
                while i < len(tokens) and tokens[i] not in "MmLlHhVvCcSsQqTtAaZz":
                    x, i = _pop_float(tokens, i)
                    y, i = _pop_float(tokens, i)
                    if cmd == "m":
                        x += cur_x
                        y += cur_y
                    cur_x, cur_y = x, y
                    ln = etree.SubElement(path_el, qn("a:lnTo"))
                    pt = etree.SubElement(ln, qn("a:pt"))
                    pt.set("x", str(_svg_to_emu(x)))
                    pt.set("y", str(_svg_to_emu(y)))

            elif cmd in ("L", "l"):
                while i < len(tokens) and tokens[i] not in "MmLlHhVvCcSsQqTtAaZz":
                    x, i = _pop_float(tokens, i)
                    y, i = _pop_float(tokens, i)
                    if cmd == "l":
                        x += cur_x
                        y += cur_y
                    cur_x, cur_y = x, y
                    ln = etree.SubElement(path_el, qn("a:lnTo"))
                    pt = etree.SubElement(ln, qn("a:pt"))
                    pt.set("x", str(_svg_to_emu(x)))
                    pt.set("y", str(_svg_to_emu(y)))

            elif cmd in ("H", "h"):
                x, i = _pop_float(tokens, i)
                if cmd == "h":
                    x += cur_x
                cur_x = x
                ln = etree.SubElement(path_el, qn("a:lnTo"))
                pt = etree.SubElement(ln, qn("a:pt"))
                pt.set("x", str(_svg_to_emu(cur_x)))
                pt.set("y", str(_svg_to_emu(cur_y)))

            elif cmd in ("V", "v"):
                y, i = _pop_float(tokens, i)
                if cmd == "v":
                    y += cur_y
                cur_y = y
                ln = etree.SubElement(path_el, qn("a:lnTo"))
                pt = etree.SubElement(ln, qn("a:pt"))
                pt.set("x", str(_svg_to_emu(cur_x)))
                pt.set("y", str(_svg_to_emu(cur_y)))

            elif cmd in ("C", "c"):
                while i < len(tokens) and tokens[i] not in "MmLlHhVvCcSsQqTtAaZz":
                    c1x, i = _pop_float(tokens, i)
                    c1y, i = _pop_float(tokens, i)
                    c2x, i = _pop_float(tokens, i)
                    c2y, i = _pop_float(tokens, i)
                    ex, i = _pop_float(tokens, i)
                    ey, i = _pop_float(tokens, i)
                    if cmd == "c":
                        c1x += cur_x; c1y += cur_y
                        c2x += cur_x; c2y += cur_y
                        ex += cur_x; ey += cur_y
                    cur_x, cur_y = ex, ey
                    cubic = etree.SubElement(path_el, qn("a:cubicBezTo"))
                    for px, py in [(c1x, c1y), (c2x, c2y), (ex, ey)]:
                        pt = etree.SubElement(cubic, qn("a:pt"))
                        pt.set("x", str(_svg_to_emu(px)))
                        pt.set("y", str(_svg_to_emu(py)))

            elif cmd in ("S", "s"):
                # Smooth cubic — reflect previous control point
                while i < len(tokens) and tokens[i] not in "MmLlHhVvCcSsQqTtAaZz":
                    c2x, i = _pop_float(tokens, i)
                    c2y, i = _pop_float(tokens, i)
                    ex, i = _pop_float(tokens, i)
                    ey, i = _pop_float(tokens, i)
                    if cmd == "s":
                        c2x += cur_x; c2y += cur_y
                        ex += cur_x; ey += cur_y
                    c1x, c1y = cur_x, cur_y  # simplified reflection
                    cur_x, cur_y = ex, ey
                    cubic = etree.SubElement(path_el, qn("a:cubicBezTo"))
                    for px, py in [(c1x, c1y), (c2x, c2y), (ex, ey)]:
                        pt = etree.SubElement(cubic, qn("a:pt"))
                        pt.set("x", str(_svg_to_emu(px)))
                        pt.set("y", str(_svg_to_emu(py)))

            elif cmd in ("A", "a"):
                while i < len(tokens) and tokens[i] not in "MmLlHhVvCcSsQqTtAaZz":
                    rx_v, i = _pop_float(tokens, i)
                    ry_v, i = _pop_float(tokens, i)
                    x_rot, i = _pop_float(tokens, i)
                    fa, i = _pop_float(tokens, i)
                    fs, i = _pop_float(tokens, i)
                    ex, i = _pop_float(tokens, i)
                    ey, i = _pop_float(tokens, i)
                    if cmd == "a":
                        ex += cur_x
                        ey += cur_y

                    if rx_v == 0 or ry_v == 0:
                        # Degenerate arc → line
                        cur_x, cur_y = ex, ey
                        ln = etree.SubElement(path_el, qn("a:lnTo"))
                        pt = etree.SubElement(ln, qn("a:pt"))
                        pt.set("x", str(_svg_to_emu(ex)))
                        pt.set("y", str(_svg_to_emu(ey)))
                        continue

                    phi = math.radians(x_rot)
                    cx, cy, adj_rx, adj_ry, theta1, dtheta = _endpoint_to_center(
                        cur_x, cur_y, ex, ey, int(fa), int(fs), abs(rx_v), abs(ry_v), phi,
                    )
                    cos_phi = math.cos(phi)
                    sin_phi = math.sin(phi)
                    beziers = _arc_to_beziers(cx, cy, adj_rx, adj_ry, theta1, dtheta, cos_phi, sin_phi)
                    for c1x, c1y, c2x, c2y, bex, bey in beziers:
                        cubic = etree.SubElement(path_el, qn("a:cubicBezTo"))
                        for px, py in [(c1x, c1y), (c2x, c2y), (bex, bey)]:
                            pt = etree.SubElement(cubic, qn("a:pt"))
                            pt.set("x", str(_svg_to_emu(px)))
                            pt.set("y", str(_svg_to_emu(py)))
                    cur_x, cur_y = ex, ey

            elif cmd in ("Z", "z"):
                etree.SubElement(path_el, qn("a:close"))
                cur_x, cur_y = start_x, start_y

    return cust_geom


# ---------------------------------------------------------------------------
# Extract <path d="..."> from SVG markup
# ---------------------------------------------------------------------------


def _extract_svg_paths(svg_markup: str) -> List[str]:
    """Extract all ``d`` attributes from ``<path>`` elements in SVG markup."""
    return re.findall(r'<path[^>]+d="([^"]+)"', svg_markup)


# ---------------------------------------------------------------------------
# Public API — render icon block to slide
# ---------------------------------------------------------------------------


def render_icon_to_slide(
    *,
    slide,
    block: FreeformBlock,
    left: int,
    top: int,
    width: int,
    height: int,
    color_rgb: RGBColor,
) -> None:
    """Add a Tabler icon as a custGeom vector shape to the slide."""
    # Resolve SVG markup
    icon_svg = str(block.icon_svg or "").strip()
    if not icon_svg:
        raw_name = str(block.icon or block.content or "").strip()
        resolved = resolve_icon_name(raw_name, fallback="sparkles")
        svg_map = load_inline_svg_map()
        icon_svg = str(svg_map.get(resolved) or "").strip()

    if not icon_svg:
        # Ultimate fallback — render a small placeholder rect
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color_rgb
        shape.line.fill.background()
        return

    path_strings = _extract_svg_paths(icon_svg)
    if not path_strings:
        return

    # Create a placeholder rect shape, then replace its geometry with custGeom
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height,
    )
    shape.fill.background()  # icons are stroke-only, no fill

    # Set line (stroke) colour and width
    shape.line.color.rgb = color_rgb
    shape.line.width = Emu(int(2.0 * _EMU_PER_SVG))  # 2 SVG units stroke

    # Replace prstGeom with custGeom
    sp = shape._element
    sp_pr = sp.find(qn("a:spPr") if sp.find(qn("a:spPr")) is not None else qn("p:spPr"))
    if sp_pr is None:
        sp_pr = sp.find(qn("p:spPr"))
    if sp_pr is None:
        logger.warning("svg_to_shape: could not find spPr element")
        return

    # Remove existing geometry
    for tag in (qn("a:prstGeom"), qn("a:custGeom")):
        existing = sp_pr.find(tag)
        if existing is not None:
            sp_pr.remove(existing)

    cust_geom = svg_paths_to_ooxml(path_strings)
    sp_pr.append(cust_geom)
