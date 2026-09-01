"""Render charts as native python-pptx chart objects for PPTX output.

Consumes the same (chart_type, payload) tuple that render_utils.resolve_chart_spec
produces, so both HTML and PPTX renderers share the same data extraction logic.
"""
from __future__ import annotations

import logging
from typing import Any, Dict, List

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Data extraction helpers (mirror chart_svg_renderer logic)
# ---------------------------------------------------------------------------


def _numeric_series(payload: Dict[str, Any]) -> List[Dict[str, Any]]:
    series_payload = list(payload.get("series") or [])
    normalised: List[Dict[str, Any]] = []
    for item in series_payload:
        if not isinstance(item, dict):
            continue
        values: List[float] = []
        for raw in list(item.get("values") or []):
            try:
                values.append(float(raw))
            except (TypeError, ValueError):
                continue
        if values:
            normalised.append({
                "name": str(item.get("name") or f"Series {len(normalised) + 1}").strip(),
                "values": values,
            })
    return normalised


def _categories(payload: Dict[str, Any], fallback_len: int) -> List[str]:
    cats = [str(item).strip() for item in list(payload.get("categories") or []) if str(item).strip()]
    if cats:
        return cats
    return [f"P{idx + 1}" for idx in range(max(1, fallback_len))]


def _pie_values(payload: Dict[str, Any]) -> List[float]:
    values: List[float] = []
    for raw in list(payload.get("values") or []):
        try:
            values.append(float(raw))
        except (TypeError, ValueError):
            continue
    if not values:
        series = _numeric_series(payload)
        if series:
            values = list(series[0].get("values") or [])
    return values or [40.0, 30.0, 20.0, 10.0]


def _hex_to_rgb(value: str) -> RGBColor:
    token = str(value or "").strip().lstrip("#")
    if len(token) == 3:
        token = "".join(ch * 2 for ch in token)
    if len(token) != 6:
        token = "2563eb"
    return RGBColor(int(token[0:2], 16), int(token[2:4], 16), int(token[4:6], 16))


def _theme_colors(theme: Dict[str, Any]) -> List[str]:
    accent = str(theme.get("accent_color") or "#2563eb")
    return [accent, "#06b6d4", "#8b5cf6", "#f59e0b", "#10b981"]


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def render_chart_to_slide(
    *,
    slide,
    chart_type: str,
    payload: Dict[str, Any],
    theme: Dict[str, Any],
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    """Add a native chart shape to the slide."""
    ctype = str(chart_type or "").strip().lower()
    if ctype == "bar":
        _render_bar(slide, payload, theme, left, top, width, height)
    elif ctype == "line":
        _render_line(slide, payload, theme, left, top, width, height)
    elif ctype == "pie":
        _render_pie(slide, payload, theme, left, top, width, height)
    else:
        logger.warning("pptx_chart: unsupported chart_type=%s", ctype)


# ---------------------------------------------------------------------------
# Chart implementations
# ---------------------------------------------------------------------------


def _render_bar(slide, payload, theme, left, top, width, height):
    series_data = _numeric_series(payload)
    if not series_data:
        series_data = [{"name": "Series", "values": [52.0, 47.0, 44.0, 39.0]}]
    cats = _categories(payload, max(len(s["values"]) for s in series_data))

    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series_data:
        vals = s["values"]
        # Pad or truncate to match category count
        padded = (vals + [0.0] * len(cats))[:len(cats)]
        chart_data.add_series(s["name"], padded)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    title_text = str(payload.get("title") or "").strip()
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text

    # Apply theme colours
    colors = _theme_colors(theme)
    plot = chart.plots[0]
    plot.gap_width = 80
    for idx, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _hex_to_rgb(colors[idx % len(colors)])


def _render_line(slide, payload, theme, left, top, width, height):
    series_data = _numeric_series(payload)
    if not series_data:
        series_data = [{"name": "Series", "values": [42.0, 55.0, 68.0, 79.0]}]
    cats = _categories(payload, max(len(s["values"]) for s in series_data))

    chart_data = CategoryChartData()
    chart_data.categories = cats
    for s in series_data:
        padded = (s["values"] + [0.0] * len(cats))[:len(cats)]
        chart_data.add_series(s["name"], padded)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    title_text = str(payload.get("title") or "").strip()
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text

    colors = _theme_colors(theme)
    plot = chart.plots[0]
    for idx, series in enumerate(plot.series):
        series.format.line.color.rgb = _hex_to_rgb(colors[idx % len(colors)])
        series.format.line.width = Pt(2.5)
        series.smooth = True


def _render_pie(slide, payload, theme, left, top, width, height):
    values = _pie_values(payload)
    cats = _categories(payload, len(values))
    if len(cats) < len(values):
        cats.extend([f"S{i + 1}" for i in range(len(cats), len(values))])

    chart_data = CategoryChartData()
    chart_data.categories = cats
    chart_data.add_series("Values", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    title_text = str(payload.get("title") or "").strip()
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text

    colors = _theme_colors(theme)
    plot = chart.plots[0]
    for idx, point in enumerate(plot.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = _hex_to_rgb(colors[idx % len(colors)])
