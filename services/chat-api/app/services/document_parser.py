from __future__ import annotations

import os
import logging
import tempfile
import mimetypes
import re
from io import BytesIO
from typing import Any, Dict, List, Optional
import csv

import httpx
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.core.config import get_settings
from app.services.document_processing_parse_client import document_processing_parse_client
from app.skills_specs.pdf.scripts.extract_pdf_text import extract_text as extract_pdf_text
from app.utils.oss_uploader import AliyunOSSUploader


logger = logging.getLogger(__name__)


class DocumentParserService:
    MAX_SPREADSHEET_ROWS_PER_SHEET = 80
    MAX_SPREADSHEET_COLS_PER_SHEET = 30
    MAX_SPREADSHEET_SHEETS = 12
    MAX_EMBEDDED_DOCX_IMAGES = 80
    MAX_PDF_TABLE_ROWS = 120
    MAX_PDF_TABLE_COLS = 24

    def __init__(self) -> None:
        self._settings = get_settings()
        self._endpoint = str(self._settings.DOC_PARSER_API_URL or "").strip()

    @staticmethod
    def _resolve_document_url(document: Dict[str, Any]) -> str:
        urls = DocumentParserService._resolve_document_urls(document)
        return urls[0] if urls else ""

    @staticmethod
    def _resolve_document_urls(document: Dict[str, Any]) -> List[str]:
        urls: List[str] = []

        def _add(value: str) -> None:
            candidate = str(value or "").strip()
            if candidate and candidate not in urls:
                urls.append(candidate)

        object_path = str(document.get("object_path") or "").strip()
        if object_path:
            try:
                uploader = AliyunOSSUploader()
                _add(str(uploader.internal_url(object_path) or "").strip())
                _add(str(uploader.sign_url(object_path) or "").strip())
            except Exception:
                pass
        for key in ("signed_url", "url"):
            value = str(document.get(key) or "").strip()
            if value:
                # Keep externally supplied OSS URLs exactly as-is first. Upstream
                # services may use a different signing version or bucket, and
                # their query string must survive unchanged.
                _add(value)
                try:
                    uploader = AliyunOSSUploader()
                    derived_object_path = uploader.object_path_from_url(value)
                    if derived_object_path:
                        signed = str(uploader.sign_url(derived_object_path) or "").strip()
                        _add(signed)
                    external_signed = str(uploader.sign_url_from_url(value) or "").strip()
                    _add(external_signed)
                except Exception:
                    pass
        return urls

    @staticmethod
    def resolve_document_object_path(document: Dict[str, Any]) -> str:
        object_path = str(document.get("object_path") or "").strip()
        if object_path:
            return object_path
        try:
            uploader = AliyunOSSUploader()
            for key in ("signed_url", "url"):
                derived = uploader.object_path_from_url(str(document.get(key) or "").strip())
                if derived:
                    return derived
        except Exception:
            return ""
        return ""

    @staticmethod
    def _join_blocks(blocks: List[str]) -> str:
        out: List[str] = []
        for block in blocks:
            text = str(block or "").strip()
            if not text:
                continue
            if out and out[-1] == text:
                continue
            out.append(text)
        return "\n\n".join(out).strip()

    @staticmethod
    def _document_ext(filename: str, source_url: str) -> str:
        token = str(filename or "").strip().lower()
        if "." in token:
            return token.rsplit(".", 1)[-1]
        url = str(source_url or "").strip().lower()
        if "." in url.split("?", 1)[0]:
            return url.split("?", 1)[0].rsplit(".", 1)[-1]
        return ""

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        rows: List[List[str]] = []
        for row in list(getattr(table, "rows", []) or []):
            values: List[str] = []
            for cell in list(getattr(row, "cells", []) or []):
                text = str(getattr(cell, "text", "") or "").replace("\n", " ").strip()
                values.append(text)
            if any(value for value in values):
                rows.append(values)
        if not rows:
            return ""
        col_count = max(len(row) for row in rows)
        normalized_rows = [row + [""] * (col_count - len(row)) for row in rows]
        header = normalized_rows[0]
        separator = ["---"] * col_count
        body = normalized_rows[1:]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(separator) + " |",
        ]
        for row in body:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines).strip()

    @staticmethod
    def _element_local_name(element: Any) -> str:
        return str(getattr(element, "tag", "") or "").split("}", 1)[-1]

    @classmethod
    def _paragraph_image_rel_ids(cls, paragraph: Any) -> List[str]:
        rel_ids: List[str] = []
        seen = set()
        root = getattr(paragraph, "_p", None)
        if root is None:
            return rel_ids
        for element in root.iter():
            if cls._element_local_name(element) != "blip":
                continue
            for key, value in dict(getattr(element, "attrib", {}) or {}).items():
                attr_name = str(key or "").split("}", 1)[-1]
                if attr_name not in {"embed", "link"}:
                    continue
                rel_id = str(value or "").strip()
                if rel_id and rel_id not in seen:
                    seen.add(rel_id)
                    rel_ids.append(rel_id)
        return rel_ids

    @staticmethod
    def _nearby_text(texts: List[str], index: int, *, before: int = 3, after: int = 2) -> Dict[str, str]:
        before_values = [str(x or "").strip() for x in texts[max(0, index - before) : index] if str(x or "").strip()]
        after_values = [str(x or "").strip() for x in texts[index + 1 : index + 1 + after] if str(x or "").strip()]
        return {
            "before_text": "\n".join(before_values[-before:]).strip(),
            "after_text": "\n".join(after_values[:after]).strip(),
            "near_text": "\n".join((before_values[-before:] + after_values[:after])[: before + after]).strip(),
        }

    @staticmethod
    def _image_extension(*, content_type: str, partname: str) -> str:
        guessed = mimetypes.guess_extension(str(content_type or "").strip()) or ""
        if guessed:
            return guessed.lstrip(".").replace("jpeg", "jpg")
        token = os.path.basename(str(partname or "").split("?", 1)[0])
        if "." in token:
            return token.rsplit(".", 1)[-1].lower()
        return "png"

    @classmethod
    def _build_docx_parse_from_bytes(
        cls,
        content: bytes,
        *,
        filename: str = "",
        upload_images: bool = False,
        user_id: str = "document-parser",
    ) -> Dict[str, Any]:
        if not content:
            return {"markdown": "", "embedded_images": []}
        doc = DocxDocument(BytesIO(content))
        paragraphs = list(getattr(doc, "paragraphs", []) or [])
        tables = list(getattr(doc, "tables", []) or [])
        paragraph_texts = [str(getattr(para, "text", "") or "").strip() for para in paragraphs]
        para_idx = 0
        table_idx = 0
        image_idx = 0
        blocks: List[str] = []
        embedded_images: List[Dict[str, Any]] = []
        uploader: Optional[AliyunOSSUploader] = None
        if upload_images:
            try:
                uploader = AliyunOSSUploader()
            except Exception:
                uploader = None

        body = getattr(getattr(doc, "element", None), "body", None)
        children = list(getattr(body, "iterchildren", lambda: [])() or [])
        for child in children:
            tag = cls._element_local_name(child)
            if tag == "p":
                if para_idx >= len(paragraphs):
                    continue
                paragraph = paragraphs[para_idx]
                text = paragraph_texts[para_idx]
                rel_ids = cls._paragraph_image_rel_ids(paragraph)
                current_para_idx = para_idx
                para_idx += 1
                if text:
                    blocks.append(text)
                if rel_ids:
                    nearby = cls._nearby_text(paragraph_texts, current_para_idx)
                    for rel_id in rel_ids:
                        if image_idx >= cls.MAX_EMBEDDED_DOCX_IMAGES:
                            break
                        part = getattr(getattr(doc, "part", None), "related_parts", {}).get(rel_id)
                        if part is None:
                            continue
                        blob = bytes(getattr(part, "blob", b"") or b"")
                        if not blob:
                            continue
                        image_idx += 1
                        content_type = str(getattr(part, "content_type", "") or "image/png").strip()
                        partname = str(getattr(part, "partname", "") or "")
                        ext = cls._image_extension(content_type=content_type, partname=partname)
                        image_filename = f"{os.path.splitext(os.path.basename(filename or 'document'))[0]}_image_{image_idx:03d}.{ext}"
                        url = ""
                        object_path = ""
                        if uploader is not None:
                            try:
                                url, object_path = uploader.upload_bytes_with_path(
                                    blob,
                                    user_id=user_id or "document-parser",
                                    file_name=image_filename,
                                    content_type=content_type,
                                )
                            except Exception:
                                url = ""
                                object_path = ""
                        caption_seed = text or nearby.get("before_text") or nearby.get("after_text") or image_filename
                        image_record = {
                            "image_id": f"embedded_image_{image_idx}",
                            "image_index": image_idx,
                            "source": "embedded_docx_image",
                            "filename": image_filename,
                            "content_type": content_type,
                            "size": len(blob),
                            "object_path": object_path,
                            "signed_url": url,
                            "url": url,
                            "paragraph_index": current_para_idx,
                            "source_order": image_idx,
                            "caption_seed": str(caption_seed or "").strip()[:240],
                            **nearby,
                        }
                        embedded_images.append(image_record)
                        context_bits = [
                            f"图片{image_idx}",
                            f"文件名：{image_filename}",
                        ]
                        near_text = str(image_record.get("near_text") or "").strip()
                        if near_text:
                            context_bits.append(f"邻近文本：{near_text[:500]}")
                        blocks.append("【内嵌图片：" + "；".join(context_bits) + "】")
                continue
            if tag == "tbl":
                if table_idx >= len(tables):
                    continue
                table_markdown = cls._table_to_markdown(tables[table_idx])
                table_idx += 1
                if table_markdown:
                    blocks.append(table_markdown)
        return {
            "markdown": cls._join_blocks(blocks),
            "embedded_images": embedded_images,
        }

    @staticmethod
    def _escape_markdown_cell(value: Any) -> str:
        text = "" if value is None else str(value)
        text = text.replace("\r", " ").replace("\n", " ").strip()
        text = text.replace("|", "\\|")
        return text

    @classmethod
    def _rows_to_markdown_table(cls, rows: List[List[Any]], *, max_rows: int, max_cols: int) -> str:
        if not rows:
            return ""
        normalized: List[List[str]] = []
        for row in rows[:max_rows]:
            cells = list(row or [])[:max_cols]
            normalized.append([cls._escape_markdown_cell(cell) for cell in cells])
        if not normalized:
            return ""
        col_count = max(len(row) for row in normalized)
        if col_count <= 0:
            return ""
        normalized = [row + [""] * (col_count - len(row)) for row in normalized]
        header = normalized[0]
        if not any(cell for cell in header):
            header = [f"Column {idx}" for idx in range(1, col_count + 1)]
            body = normalized
        else:
            body = normalized[1:]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * col_count) + " |",
        ]
        for row in body:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines).strip()

    @staticmethod
    def _bbox_contains(outer: Any, inner: Any, *, tolerance: float = 1.5) -> bool:
        try:
            ox0, oy0, ox1, oy1 = [float(x) for x in list(outer or [])[:4]]
            ix0, iy0, ix1, iy1 = [float(x) for x in list(inner or [])[:4]]
        except Exception:
            return False
        return (
            ix0 >= ox0 - tolerance
            and iy0 >= oy0 - tolerance
            and ix1 <= ox1 + tolerance
            and iy1 <= oy1 + tolerance
        )

    @classmethod
    def _pdf_words_to_paragraphs(cls, words: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        normalized: List[Dict[str, Any]] = []
        for word in words:
            if not isinstance(word, dict):
                continue
            text = str(word.get("text") or "").strip()
            if not text:
                continue
            try:
                normalized.append(
                    {
                        "text": text,
                        "x0": float(word.get("x0") or 0),
                        "x1": float(word.get("x1") or 0),
                        "top": float(word.get("top") or 0),
                        "bottom": float(word.get("bottom") or 0),
                    }
                )
            except Exception:
                continue
        if not normalized:
            return []

        normalized.sort(key=lambda item: (round(float(item["top"]) / 3.0), float(item["x0"])))
        lines: List[Dict[str, Any]] = []
        for word in normalized:
            matched = None
            for line in reversed(lines[-5:]):
                if abs(float(line["top"]) - float(word["top"])) <= 3.0:
                    matched = line
                    break
            if matched is None:
                matched = {
                    "words": [],
                    "x0": float(word["x0"]),
                    "x1": float(word["x1"]),
                    "top": float(word["top"]),
                    "bottom": float(word["bottom"]),
                }
                lines.append(matched)
            matched["words"].append(word)
            matched["x0"] = min(float(matched["x0"]), float(word["x0"]))
            matched["x1"] = max(float(matched["x1"]), float(word["x1"]))
            matched["top"] = min(float(matched["top"]), float(word["top"]))
            matched["bottom"] = max(float(matched["bottom"]), float(word["bottom"]))

        cleaned_lines: List[Dict[str, Any]] = []
        for line in lines:
            line_words = sorted(list(line.get("words") or []), key=lambda item: float(item["x0"]))
            text = " ".join(str(item.get("text") or "").strip() for item in line_words if str(item.get("text") or "").strip())
            text = re.sub(r"\s+([,.;:%)\]])", r"\1", text)
            text = re.sub(r"([(\[])\s+", r"\1", text).strip()
            if not text:
                continue
            cleaned_lines.append(
                {
                    "text": text,
                    "bbox": [float(line["x0"]), float(line["top"]), float(line["x1"]), float(line["bottom"])],
                }
            )

        paragraphs: List[Dict[str, Any]] = []
        current: List[Dict[str, Any]] = []

        def _flush() -> None:
            nonlocal current
            if not current:
                return
            text = " ".join(str(item.get("text") or "").strip() for item in current if str(item.get("text") or "").strip())
            text = re.sub(r"\s+", " ", text).strip()
            if text:
                x0 = min(float(item["bbox"][0]) for item in current)
                top = min(float(item["bbox"][1]) for item in current)
                x1 = max(float(item["bbox"][2]) for item in current)
                bottom = max(float(item["bbox"][3]) for item in current)
                paragraphs.append({"text": text, "bbox": [x0, top, x1, bottom]})
            current = []

        for line in cleaned_lines:
            if not current:
                current.append(line)
                continue
            prev = current[-1]
            prev_bbox = list(prev.get("bbox") or [0, 0, 0, 0])
            curr_bbox = list(line.get("bbox") or [0, 0, 0, 0])
            vertical_gap = float(curr_bbox[1]) - float(prev_bbox[3])
            left_delta = abs(float(curr_bbox[0]) - float(prev_bbox[0]))
            prev_text = str(prev.get("text") or "").strip()
            sentence_ended = bool(re.search(r"[。！？.!?;；:]$|[。！？.!?]\s*[\"'”’)]?$", prev_text))
            starts_list = bool(re.match(r"^(\d+[\.)]|[A-Z][\.)]|[-•])\s+", str(line.get("text") or "")))
            if vertical_gap > 12 or left_delta > 36 or starts_list or (vertical_gap > 6 and sentence_ended):
                _flush()
            current.append(line)
        _flush()
        return paragraphs

    @classmethod
    def _score_pdf_parse_quality(cls, markdown: str, *, table_count: int, page_count: int) -> Dict[str, Any]:
        text = str(markdown or "")
        if not text:
            return {"level": "failed", "reason": "empty_markdown", "score": 0.0}
        compact = re.sub(r"\s+", "", text)
        if not compact:
            return {"level": "failed", "reason": "empty_markdown", "score": 0.0}
        suspicious = len(re.findall(r"[^\w\s\u4e00-\u9fff.,;:!?%$€£¥()\[\]{}<>/'\"“”‘’\-–—+*=#&@|]", compact))
        suspicious_ratio = suspicious / max(len(compact), 1)
        very_long_lines = sum(1 for line in text.splitlines() if len(line.strip()) > 240)
        score = 1.0
        reasons: List[str] = []
        if suspicious_ratio > 0.08:
            score -= 0.35
            reasons.append("high_suspicious_char_ratio")
        if very_long_lines > max(3, page_count):
            score -= 0.25
            reasons.append("many_overlong_lines")
        if table_count == 0 and re.search(r"\b(Revenue|EBITA|GAAP|RMB|cash flow|balance sheet)\b", text, re.I):
            score -= 0.15
            reasons.append("financial_text_without_tables")
        level = "high" if score >= 0.8 else "medium" if score >= 0.55 else "low"
        return {
            "level": level,
            "score": round(max(score, 0.0), 3),
            "suspicious_char_ratio": round(suspicious_ratio, 4),
            "overlong_lines": very_long_lines,
            "reasons": reasons,
        }

    @classmethod
    def _build_markdown_from_local_xlsx_bytes(cls, content: bytes, *, filename: str = "") -> str:
        if not content:
            return ""
        import openpyxl

        wb = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
        blocks: List[str] = []
        title = str(filename or "spreadsheet").strip()
        if title:
            blocks.append(f"Source spreadsheet: {title}")
        for sheet_idx, sheet_name in enumerate(list(wb.sheetnames or [])[: cls.MAX_SPREADSHEET_SHEETS], start=1):
            ws = wb[sheet_name]
            rows: List[List[Any]] = []
            non_empty_rows = 0
            for row in ws.iter_rows(values_only=True):
                values = list(row or [])
                if not any(cell is not None and str(cell).strip() for cell in values):
                    continue
                non_empty_rows += 1
                if len(rows) < cls.MAX_SPREADSHEET_ROWS_PER_SHEET:
                    rows.append(values[: cls.MAX_SPREADSHEET_COLS_PER_SHEET])
            if not rows:
                continue
            used_rows = int(getattr(ws, "max_row", 0) or 0)
            used_cols = int(getattr(ws, "max_column", 0) or 0)
            block_lines = [
                f"## Sheet: {sheet_name}",
                "",
                f"- Sheet index: {sheet_idx}",
                f"- Used range: {used_rows} rows x {used_cols} columns",
                f"- Non-empty rows detected: {non_empty_rows}",
            ]
            if non_empty_rows > cls.MAX_SPREADSHEET_ROWS_PER_SHEET or used_cols > cls.MAX_SPREADSHEET_COLS_PER_SHEET:
                block_lines.append(
                    f"- Preview truncated to first {cls.MAX_SPREADSHEET_ROWS_PER_SHEET} non-empty rows "
                    f"and {cls.MAX_SPREADSHEET_COLS_PER_SHEET} columns."
                )
            table = cls._rows_to_markdown_table(
                rows,
                max_rows=cls.MAX_SPREADSHEET_ROWS_PER_SHEET,
                max_cols=cls.MAX_SPREADSHEET_COLS_PER_SHEET,
            )
            if table:
                block_lines.extend(["", table])
            blocks.append("\n".join(block_lines).strip())
        if len(list(wb.sheetnames or [])) > cls.MAX_SPREADSHEET_SHEETS:
            blocks.append(f"_Additional sheets omitted: {len(wb.sheetnames) - cls.MAX_SPREADSHEET_SHEETS}_")
        try:
            wb.close()
        except Exception:
            pass
        return cls._join_blocks(blocks)

    @classmethod
    def _build_markdown_from_local_delimited_bytes(cls, content: bytes, *, filename: str = "", delimiter: str = ",") -> str:
        if not content:
            return ""
        text = content.decode("utf-8-sig", errors="replace")
        reader = csv.reader(text.splitlines(), delimiter=delimiter)
        rows: List[List[Any]] = []
        non_empty_rows = 0
        max_cols_seen = 0
        for row in reader:
            values = list(row or [])
            if not any(str(cell or "").strip() for cell in values):
                continue
            non_empty_rows += 1
            max_cols_seen = max(max_cols_seen, len(values))
            if len(rows) < cls.MAX_SPREADSHEET_ROWS_PER_SHEET:
                rows.append(values[: cls.MAX_SPREADSHEET_COLS_PER_SHEET])
        table = cls._rows_to_markdown_table(
            rows,
            max_rows=cls.MAX_SPREADSHEET_ROWS_PER_SHEET,
            max_cols=cls.MAX_SPREADSHEET_COLS_PER_SHEET,
        )
        if not table:
            return ""
        blocks = [
            f"Source spreadsheet: {str(filename or 'delimited file').strip()}",
            "## Sheet: Data",
            "",
            f"- Non-empty rows detected: {non_empty_rows}",
            f"- Columns detected: {max_cols_seen}",
        ]
        if non_empty_rows > cls.MAX_SPREADSHEET_ROWS_PER_SHEET or max_cols_seen > cls.MAX_SPREADSHEET_COLS_PER_SHEET:
            blocks.append(
                f"- Preview truncated to first {cls.MAX_SPREADSHEET_ROWS_PER_SHEET} non-empty rows "
                f"and {cls.MAX_SPREADSHEET_COLS_PER_SHEET} columns."
            )
        blocks.extend(["", table])
        return "\n".join(blocks).strip()

    @classmethod
    def _build_markdown_from_local_xls_bytes(cls, content: bytes, *, filename: str = "") -> str:
        if not content:
            return ""
        try:
            import xlrd
        except ImportError as exc:
            raise RuntimeError("xlrd is required to parse legacy .xls files") from exc

        book = xlrd.open_workbook(file_contents=content)
        blocks: List[str] = []
        title = str(filename or "spreadsheet").strip()
        if title:
            blocks.append(f"Source spreadsheet: {title}")
        for sheet_idx in range(min(book.nsheets, cls.MAX_SPREADSHEET_SHEETS)):
            sheet = book.sheet_by_index(sheet_idx)
            rows: List[List[Any]] = []
            non_empty_rows = 0
            for ridx in range(sheet.nrows):
                values = sheet.row_values(ridx)
                if not any(str(cell or "").strip() for cell in values):
                    continue
                non_empty_rows += 1
                if len(rows) < cls.MAX_SPREADSHEET_ROWS_PER_SHEET:
                    rows.append(values[: cls.MAX_SPREADSHEET_COLS_PER_SHEET])
            if not rows:
                continue
            block_lines = [
                f"## Sheet: {sheet.name}",
                "",
                f"- Sheet index: {sheet_idx + 1}",
                f"- Used range: {sheet.nrows} rows x {sheet.ncols} columns",
                f"- Non-empty rows detected: {non_empty_rows}",
            ]
            if non_empty_rows > cls.MAX_SPREADSHEET_ROWS_PER_SHEET or sheet.ncols > cls.MAX_SPREADSHEET_COLS_PER_SHEET:
                block_lines.append(
                    f"- Preview truncated to first {cls.MAX_SPREADSHEET_ROWS_PER_SHEET} non-empty rows "
                    f"and {cls.MAX_SPREADSHEET_COLS_PER_SHEET} columns."
                )
            table = cls._rows_to_markdown_table(
                rows,
                max_rows=cls.MAX_SPREADSHEET_ROWS_PER_SHEET,
                max_cols=cls.MAX_SPREADSHEET_COLS_PER_SHEET,
            )
            if table:
                block_lines.extend(["", table])
            blocks.append("\n".join(block_lines).strip())
        if book.nsheets > cls.MAX_SPREADSHEET_SHEETS:
            blocks.append(f"_Additional sheets omitted: {book.nsheets - cls.MAX_SPREADSHEET_SHEETS}_")
        return cls._join_blocks(blocks)

    @classmethod
    def _shape_blocks(cls, shape: Any) -> List[str]:
        blocks: List[str] = []
        try:
            shape_type = getattr(shape, "shape_type", None)
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                children = sorted(
                    list(getattr(shape, "shapes", []) or []),
                    key=lambda item: (
                        int(getattr(item, "top", 0) or 0),
                        int(getattr(item, "left", 0) or 0),
                    ),
                )
                for child in children:
                    blocks.extend(cls._shape_blocks(child))
                return blocks

            if bool(getattr(shape, "has_table", False)):
                table_markdown = cls._table_to_markdown(getattr(shape, "table", None))
                if table_markdown:
                    blocks.append(table_markdown)
                return blocks

            text = ""
            if hasattr(shape, "text"):
                text = str(getattr(shape, "text", "") or "").strip()
            if text:
                blocks.append(text)
        except Exception:
            return blocks
        return blocks

    @classmethod
    def _build_markdown_from_local_pptx_bytes(cls, content: bytes) -> str:
        presentation = Presentation(BytesIO(content))
        slides: List[str] = []
        for idx, slide in enumerate(list(presentation.slides or []), start=1):
            ordered_shapes = sorted(
                list(getattr(slide, "shapes", []) or []),
                key=lambda item: (
                    int(getattr(item, "top", 0) or 0),
                    int(getattr(item, "left", 0) or 0),
                ),
            )
            page_blocks: List[str] = []
            for shape in ordered_shapes:
                page_blocks.extend(cls._shape_blocks(shape))
            page_body = cls._join_blocks(page_blocks)
            if page_body:
                slides.append(f"## Page {idx}\n\n{page_body}")
        return cls._join_blocks(slides)

    @classmethod
    def _build_pdf_parse_from_bytes(cls, content: bytes, *, filename: str = "") -> Dict[str, Any]:
        if not content:
            return {
                "markdown": "",
                "structured_content": {"pages": [], "tables": [], "images": []},
                "embedded_images": [],
                "parse_quality": {"level": "failed", "reason": "empty_content", "score": 0.0},
                "parser": "local_pdf_structured",
            }
        tmp_path = ""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            pages: List[Dict[str, Any]] = []
            all_tables: List[Dict[str, Any]] = []
            all_images: List[Dict[str, Any]] = []
            markdown_pages: List[str] = []
            parser_name = "local_pdf_structured"

            try:
                import pdfplumber
            except ImportError:
                pdfplumber = None

            if pdfplumber is not None:
                try:
                    with pdfplumber.open(tmp_path) as pdf:
                        for page_index, page in enumerate(list(pdf.pages or []), start=1):
                            page_tables: List[Dict[str, Any]] = []
                            table_bboxes: List[List[float]] = []
                            try:
                                found_tables = list(page.find_tables() or [])
                            except Exception:
                                found_tables = []
                            for table_index, table in enumerate(found_tables, start=1):
                                try:
                                    rows = table.extract() or []
                                except Exception:
                                    rows = []
                                normalized_rows = [
                                    [str(cell or "").strip() for cell in list(row or [])]
                                    for row in rows
                                    if any(str(cell or "").strip() for cell in list(row or []))
                                ]
                                if not normalized_rows:
                                    continue
                                bbox = [float(x) for x in list(getattr(table, "bbox", []) or [])[:4]]
                                if len(bbox) == 4:
                                    table_bboxes.append(bbox)
                                table_markdown = cls._rows_to_markdown_table(
                                    normalized_rows,
                                    max_rows=cls.MAX_PDF_TABLE_ROWS,
                                    max_cols=cls.MAX_PDF_TABLE_COLS,
                                )
                                if not table_markdown:
                                    continue
                                table_record = {
                                    "table_id": f"pdf_page_{page_index}_table_{table_index}",
                                    "page": page_index,
                                    "bbox": bbox,
                                    "rows": normalized_rows[: cls.MAX_PDF_TABLE_ROWS],
                                    "row_count": len(normalized_rows),
                                    "col_count": max(len(row) for row in normalized_rows) if normalized_rows else 0,
                                    "table_markdown": table_markdown,
                                }
                                page_tables.append(table_record)
                                all_tables.append(table_record)

                            try:
                                words = list(
                                    page.extract_words(
                                        x_tolerance=2,
                                        y_tolerance=3,
                                        keep_blank_chars=False,
                                        use_text_flow=False,
                                    )
                                    or []
                                )
                            except Exception:
                                words = []
                            filtered_words = []
                            for word in words:
                                try:
                                    word_bbox = [
                                        float(word.get("x0") or 0),
                                        float(word.get("top") or 0),
                                        float(word.get("x1") or 0),
                                        float(word.get("bottom") or 0),
                                    ]
                                except Exception:
                                    continue
                                if any(cls._bbox_contains(table_bbox, word_bbox) for table_bbox in table_bboxes):
                                    continue
                                filtered_words.append(word)
                            paragraphs = cls._pdf_words_to_paragraphs(filtered_words)

                            page_images: List[Dict[str, Any]] = []
                            for image_index, image in enumerate(list(getattr(page, "images", []) or []), start=1):
                                try:
                                    bbox = [
                                        float(image.get("x0") or 0),
                                        float(image.get("top") or 0),
                                        float(image.get("x1") or 0),
                                        float(image.get("bottom") or 0),
                                    ]
                                    width = max(0.0, bbox[2] - bbox[0])
                                    height = max(0.0, bbox[3] - bbox[1])
                                except Exception:
                                    continue
                                if width * height < 400 or min(width, height) < 12:
                                    continue
                                image_record = {
                                    "image_id": f"pdf_page_{page_index}_image_{image_index}",
                                    "image_index": image_index,
                                    "source": "embedded_pdf_image",
                                    "filename": f"{os.path.splitext(os.path.basename(filename or 'document'))[0]}_page_{page_index}_image_{image_index}",
                                    "page": page_index,
                                    "bbox": bbox,
                                    "width": round(width, 2),
                                    "height": round(height, 2),
                                    "object_path": "",
                                    "signed_url": "",
                                    "url": "",
                                    "caption_seed": f"PDF page {page_index} image {image_index}",
                                }
                                page_images.append(image_record)
                                all_images.append(image_record)

                            page_elements: List[Dict[str, Any]] = []
                            for para_index, paragraph in enumerate(paragraphs, start=1):
                                page_elements.append(
                                    {
                                        "kind": "text",
                                        "order": len(page_elements) + 1,
                                        "bbox": paragraph.get("bbox") or [],
                                        "text": str(paragraph.get("text") or "").strip(),
                                        "paragraph_index": para_index,
                                    }
                                )
                            for table in page_tables:
                                page_elements.append(
                                    {
                                        "kind": "table",
                                        "order": len(page_elements) + 1,
                                        "bbox": table.get("bbox") or [],
                                        "table_id": table.get("table_id"),
                                        "text": str(table.get("table_markdown") or "").strip(),
                                    }
                                )
                            for image in page_images:
                                page_elements.append(
                                    {
                                        "kind": "image",
                                        "order": len(page_elements) + 1,
                                        "bbox": image.get("bbox") or [],
                                        "image_id": image.get("image_id"),
                                        "text": f"【内嵌图片：第 {page_index} 页，图片 {image.get('image_index')}】",
                                    }
                                )
                            page_elements.sort(
                                key=lambda item: (
                                    float((list(item.get("bbox") or [0, 0, 0, 0]) + [0, 0, 0, 0])[1] or 0),
                                    float((list(item.get("bbox") or [0, 0, 0, 0]) + [0, 0, 0, 0])[0] or 0),
                                )
                            )
                            body_blocks = [str(item.get("text") or "").strip() for item in page_elements if str(item.get("text") or "").strip()]
                            page_markdown = cls._join_blocks(body_blocks)
                            if page_markdown:
                                markdown_pages.append(f"## Page {page_index}\n\n{page_markdown}")
                            pages.append(
                                {
                                    "page": page_index,
                                    "width": float(getattr(page, "width", 0) or 0),
                                    "height": float(getattr(page, "height", 0) or 0),
                                    "blocks": page_elements,
                                    "paragraphs": paragraphs,
                                    "tables": page_tables,
                                    "images": page_images,
                                }
                            )
                except Exception as exc:
                    logger.warning("local structured PDF parse failed, falling back to pypdf: %s", exc)
                    pages = []
                    all_tables = []
                    all_images = []
                    markdown_pages = []

            markdown = cls._join_blocks(markdown_pages)
            if not markdown:
                parser_name = "local_pdf_text_fallback"
                try:
                    from pypdf import PdfReader

                    reader = PdfReader(tmp_path)
                    fallback_pages: List[str] = []
                    pages = []
                    for page_index, page in enumerate(list(reader.pages or []), start=1):
                        text = str(page.extract_text() or "").strip()
                        if not text:
                            continue
                        fallback_pages.append(f"## Page {page_index}\n\n{text}")
                        pages.append(
                            {
                                "page": page_index,
                                "blocks": [{"kind": "text", "order": 1, "bbox": [], "text": text}],
                                "paragraphs": [{"text": text, "bbox": []}],
                                "tables": [],
                                "images": [],
                            }
                        )
                    markdown = cls._join_blocks(fallback_pages)
                except Exception:
                    text = str(extract_pdf_text(tmp_path, max_pages=None, layout=True) or "").strip()
                    if not text:
                        text = str(extract_pdf_text(tmp_path, max_pages=None, layout=False) or "").strip()
                    markdown = text
                    pages = [{"page": 1, "blocks": [{"kind": "text", "order": 1, "bbox": [], "text": text}], "paragraphs": [{"text": text, "bbox": []}], "tables": [], "images": []}] if text else []

            parse_quality = cls._score_pdf_parse_quality(
                markdown,
                table_count=len(all_tables),
                page_count=len(pages),
            )
            return {
                "markdown": markdown,
                "structured_content": {
                    "pages": pages,
                    "tables": all_tables,
                    "images": all_images,
                    "source": parser_name,
                },
                "embedded_images": all_images,
                "parse_quality": parse_quality,
                "parser": parser_name,
            }
        finally:
            if tmp_path:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass

    @classmethod
    def _build_markdown_from_local_pdf_bytes(cls, content: bytes) -> str:
        return str(cls._build_pdf_parse_from_bytes(content).get("markdown") or "").strip()

    @classmethod
    def _build_markdown_from_local_docx_bytes(cls, content: bytes) -> str:
        return str(cls._build_docx_parse_from_bytes(content).get("markdown") or "").strip()

    async def _local_parse_fallback(
        self,
        *,
        document: Dict[str, Any],
        source_url: str,
        filename: str,
        timeout_seconds: float,
    ) -> Dict[str, Any]:
        ext = self._document_ext(filename, source_url)
        if ext not in {"pdf", "pptx", "docx", "xlsx", "xlsm", "csv", "tsv", "xls"}:
            return {
                "ok": False,
                "error": "local_fallback_unsupported",
                "markdown": "",
                "source_url": source_url,
                "filename": filename,
                "parser": "local_fallback",
            }
        try:
            parsed_url = httpx.URL(str(source_url or ""))
            if not parsed_url.scheme or parsed_url.scheme not in {"http", "https"}:
                return {
                    "ok": False,
                    "error": "local_fetch_invalid_url",
                    "markdown": "",
                    "source_url": source_url,
                    "filename": filename,
                    "parser": f"local_{ext or 'document'}",
                }
            async with httpx.AsyncClient(timeout=timeout_seconds, follow_redirects=True) as client:
                resp = await client.get(source_url)
            if resp.status_code >= 400:
                return {
                    "ok": False,
                    "error": f"local_fetch_http_{resp.status_code}",
                    "markdown": "",
                    "source_url": source_url,
                    "filename": filename,
                    "parser": f"local_{ext or 'document'}",
                }
            content = bytes(resp.content or b"")
            if ext == "pdf":
                pdf_result = self._build_pdf_parse_from_bytes(content, filename=filename)
                markdown = str(pdf_result.get("markdown") or "").strip()
                embedded_images = [dict(x) for x in list(pdf_result.get("embedded_images") or []) if isinstance(x, dict)]
                structured_content = dict(pdf_result.get("structured_content") or {})
                parse_quality = dict(pdf_result.get("parse_quality") or {})
                parser = str(pdf_result.get("parser") or f"local_{ext or 'document'}").strip()
            elif ext == "docx":
                docx_result = self._build_docx_parse_from_bytes(
                    content,
                    filename=filename,
                    upload_images=True,
                    user_id=str(document.get("user_id") or "document-parser").strip() or "document-parser",
                )
                markdown = str(docx_result.get("markdown") or "").strip()
                embedded_images = [dict(x) for x in list(docx_result.get("embedded_images") or []) if isinstance(x, dict)]
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            elif ext in {"xlsx", "xlsm"}:
                markdown = self._build_markdown_from_local_xlsx_bytes(content, filename=filename)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            elif ext == "csv":
                markdown = self._build_markdown_from_local_delimited_bytes(content, filename=filename, delimiter=",")
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            elif ext == "tsv":
                markdown = self._build_markdown_from_local_delimited_bytes(content, filename=filename, delimiter="\t")
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            elif ext == "xls":
                markdown = self._build_markdown_from_local_xls_bytes(content, filename=filename)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            else:
                markdown = self._build_markdown_from_local_pptx_bytes(content)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}"
            if self._settings.DEBUG:
                logger.debug(
                    "document local fallback parsed",
                    extra={"event": "document_parser.local_fallback_parsed", "file_name": filename, "ext": ext, "markdown_chars": len(markdown)},
                )
            return {
                "ok": bool(markdown),
                "error": "" if markdown else "local_missing_markdown",
                "markdown": markdown,
                "embedded_images": embedded_images,
                "structured_content": structured_content,
                "parse_quality": parse_quality,
                "source_url": source_url,
                "filename": filename,
                "parser": parser,
            }
        except Exception as exc:
            if self._settings.DEBUG:
                logger.warning("document local fallback failed", extra={"event": "document_parser.local_fallback_failed", "file_name": filename, "ext": ext, "error": str(exc)})
            return {
                "ok": False,
                "error": str(exc),
                "markdown": "",
                "source_url": source_url,
                "filename": filename,
                "parser": f"local_{ext or 'document'}",
            }

    async def _local_parse_object_path(
        self,
        *,
        object_path: str,
        filename: str,
    ) -> Dict[str, Any]:
        path = str(object_path or "").strip()
        ext = self._document_ext(filename, path)
        if self._settings.STORAGE_BACKEND.strip().lower() != "local" or not path:
            return {}
        if ext not in {"pdf", "pptx", "docx", "xlsx", "xlsm", "csv", "tsv", "xls"}:
            return {}
        try:
            uploader = AliyunOSSUploader()
            content = uploader.read_bytes(path)
            if ext == "pdf":
                pdf_result = self._build_pdf_parse_from_bytes(content, filename=filename)
                markdown = str(pdf_result.get("markdown") or "").strip()
                embedded_images = [dict(x) for x in list(pdf_result.get("embedded_images") or []) if isinstance(x, dict)]
                structured_content = dict(pdf_result.get("structured_content") or {})
                parse_quality = dict(pdf_result.get("parse_quality") or {})
                parser = str(pdf_result.get("parser") or f"local_{ext or 'document'}_object").strip()
            elif ext == "docx":
                docx_result = self._build_docx_parse_from_bytes(
                    content,
                    filename=filename,
                    upload_images=True,
                    user_id="document-parser",
                )
                markdown = str(docx_result.get("markdown") or "").strip()
                embedded_images = [dict(x) for x in list(docx_result.get("embedded_images") or []) if isinstance(x, dict)]
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            elif ext in {"xlsx", "xlsm"}:
                markdown = self._build_markdown_from_local_xlsx_bytes(content, filename=filename)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            elif ext == "csv":
                markdown = self._build_markdown_from_local_delimited_bytes(content, filename=filename, delimiter=",")
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            elif ext == "tsv":
                markdown = self._build_markdown_from_local_delimited_bytes(content, filename=filename, delimiter="\t")
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            elif ext == "xls":
                markdown = self._build_markdown_from_local_xls_bytes(content, filename=filename)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            else:
                markdown = self._build_markdown_from_local_pptx_bytes(content)
                embedded_images = []
                structured_content = {}
                parse_quality = {}
                parser = f"local_{ext or 'document'}_object"
            return {
                "ok": bool(markdown),
                "error": "" if markdown else "local_missing_markdown",
                "markdown": markdown,
                "embedded_images": embedded_images,
                "structured_content": structured_content,
                "parse_quality": parse_quality,
                "source_url": f"object://{path}",
                "filename": filename,
                "parser": parser,
            }
        except Exception as exc:
            if self._settings.DEBUG:
                logger.warning(
                    "document local object fallback failed",
                    extra={"event": "document_parser.local_object_fallback_failed", "file_name": filename, "ext": ext, "error": str(exc)},
                )
            return {
                "ok": False,
                "error": str(exc),
                "markdown": "",
                "source_url": f"object://{path}",
                "filename": filename,
                "parser": f"local_{ext or 'document'}_object",
            }

    @classmethod
    def _build_markdown_from_original_content(cls, data: Dict[str, Any]) -> str:
        original = data.get("original_parsed_content") if isinstance(data.get("original_parsed_content"), dict) else {}
        chunks = list(original.get("pageChunkData") or []) if isinstance(original, dict) else []
        pages: List[str] = []
        for item in chunks:
            if not isinstance(item, dict):
                continue
            page_no = int(item.get("page_number") or 0) if str(item.get("page_number") or "").strip() else 0
            page_blocks: List[str] = []

            for text in list(item.get("page_content") or []):
                value = str(text or "").strip()
                if value:
                    page_blocks.append(value)

            for table in list(item.get("page_tables") or []):
                if not isinstance(table, dict):
                    continue
                table_markdown = str(table.get("table_markdown") or "").strip()
                if table_markdown:
                    page_blocks.append(table_markdown)

            page_body = cls._join_blocks(page_blocks)
            if not page_body:
                continue
            if page_no > 0:
                pages.append(f"## Page {page_no}\n\n{page_body}")
            else:
                pages.append(page_body)
        return cls._join_blocks(pages)

    async def parse_document(
        self,
        *,
        document: Dict[str, Any],
        enable_ocr: bool = False,
        timeout_seconds: float = 90.0,
    ) -> Dict[str, Any]:
        source_urls = self._resolve_document_urls(document)
        source_url = source_urls[0] if source_urls else ""
        filename = str(document.get("filename") or "").strip()
        if self._settings.DEBUG:
            logger.debug(
                "document parse started",
                extra={"event": "document_parser.parse_started", "file_name": filename, "has_object_path": bool(str(document.get("object_path") or "").strip()), "has_url": bool(source_url), "enable_ocr": bool(enable_ocr)},
            )
        if not source_url:
            if self._settings.DEBUG:
                logger.warning("document missing source URL", extra={"event": "document_parser.missing_source_url", "file_name": filename})
            return {
                "ok": False,
                "error": "missing_document_url",
                "markdown": "",
                "source_url": "",
                "filename": filename,
                "parser": "link_markdown",
            }
        ext = self._document_ext(filename, source_url)
        local_supported_exts = {"pdf", "pptx", "docx", "xlsx", "xlsm", "csv", "tsv", "xls"}

        async def _try_local_parse_candidates() -> Dict[str, Any]:
            object_result = await self._local_parse_object_path(
                object_path=self.resolve_document_object_path(document),
                filename=filename,
            )
            if object_result.get("ok"):
                return object_result
            last: Dict[str, Any] = {}
            if object_result:
                last = object_result
            for candidate_url in source_urls:
                local_result = await self._local_parse_fallback(
                    document=document,
                    source_url=candidate_url,
                    filename=filename,
                    timeout_seconds=timeout_seconds,
                )
                last = local_result
                if local_result.get("ok"):
                    return local_result
            return last

        if ext == "pdf":
            remote = await document_processing_parse_client.parse_markdown(document, timeout_seconds=timeout_seconds)
            if remote.get("ok"):
                return remote

        if ext in {"md", "markdown", "txt"}:
            try:
                if self._settings.DEBUG:
                    logger.debug("document direct read started", extra={"event": "document_parser.direct_read_started", "file_name": filename, "ext": ext})
                async with httpx.AsyncClient(timeout=timeout_seconds, follow_redirects=True) as client:
                    resp = await client.get(source_url)
                if resp.status_code < 400:
                    markdown = str(resp.text or "").strip()
                    if self._settings.DEBUG:
                        logger.debug("document direct read succeeded", extra={"event": "document_parser.direct_read_succeeded", "file_name": filename, "chars": len(markdown)})
                    return {
                        "ok": bool(markdown),
                        "error": "" if markdown else "empty_content",
                        "markdown": markdown,
                        "source_url": source_url,
                        "filename": filename,
                        "parser": "direct_read",
                    }
                else:
                    if self._settings.DEBUG:
                        logger.warning("document direct read fetch failed", extra={"event": "document_parser.direct_read_fetch_failed", "file_name": filename, "status": resp.status_code})
            except Exception as exc:
                if self._settings.DEBUG:
                    logger.warning("document direct read failed", extra={"event": "document_parser.direct_read_failed", "file_name": filename, "error": str(exc)})

        if ext in {"pdf", "xlsx", "xlsm", "csv", "tsv", "xls"}:
            local = await _try_local_parse_candidates()
            if local.get("ok") or not self._endpoint:
                return local

        if ext == "docx":
            local = await _try_local_parse_candidates()
            if local.get("ok"):
                return local

        if not self._endpoint:
            if self._settings.DEBUG:
                logger.warning("document parser endpoint missing", extra={"event": "document_parser.endpoint_missing", "file_name": filename})
            if ext in local_supported_exts:
                local = await _try_local_parse_candidates()
                if local.get("ok"):
                    return local
            return {
                "ok": False,
                "error": "missing_DOC_PARSER_API_URL",
                "markdown": "",
                "source_url": source_url,
                "filename": filename,
                "parser": "link_markdown",
            }

        payload = {"url": source_url, "enable_ocr": bool(enable_ocr)}
        try:
            async with httpx.AsyncClient(timeout=timeout_seconds, follow_redirects=True) as client:
                resp = await client.post(self._endpoint, json=payload)
            raw_text = str(resp.text or "")
            if self._settings.DEBUG:
                logger.debug("document parser response received", extra={"event": "document_parser.response", "file_name": filename, "status": resp.status_code, "body_chars": len(raw_text)})
            if resp.status_code >= 400:
                local = await _try_local_parse_candidates()
                if local.get("ok"):
                    return local
                return {
                    "ok": False,
                    "error": f"http_{resp.status_code}",
                    "markdown": "",
                    "source_url": source_url,
                    "filename": filename,
                    "parser": "link_markdown",
                    "raw": raw_text[:1200],
                }
            data = resp.json() if raw_text else {}
        except Exception as exc:
            if self._settings.DEBUG:
                logger.warning("document parser request failed", extra={"event": "document_parser.request_failed", "file_name": filename, "error": str(exc)})
            local = await _try_local_parse_candidates()
            if local.get("ok"):
                return local
            return {
                "ok": False,
                "error": str(exc),
                "markdown": "",
                "source_url": source_url,
                "filename": filename,
                "parser": "link_markdown",
            }

        markdown = str((data or {}).get("markdown") or "").strip()
        if not markdown and isinstance(data, dict):
            markdown = self._build_markdown_from_original_content(data)
        if not markdown:
            local = await _try_local_parse_candidates()
            if local.get("ok"):
                return {
                    **local,
                    "raw_response": data if isinstance(data, dict) else {},
                }
        if self._settings.DEBUG:
            logger.debug(
                "document parsed",
                extra={"event": "document_parser.parsed", "file_name": filename, "ok": bool(markdown), "markdown_chars": len(markdown), "keys": list((data or {}).keys())[:10] if isinstance(data, dict) else []},
            )
        return {
            "ok": bool(markdown),
            "error": "" if markdown else "missing_markdown",
            "markdown": markdown,
            "source_url": source_url,
            "filename": filename,
            "parser": "link_markdown",
            "raw_response": data if isinstance(data, dict) else {},
        }


document_parser_service = DocumentParserService()
