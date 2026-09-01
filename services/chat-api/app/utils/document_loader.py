from app.infrastructure.observability.config import log_print
import io
import os
from typing import Optional

def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """
    Extract text content from various file formats.
    Supported: .docx, .pdf, .pptx, .xlsx, .txt, .md
    """
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == '.docx':
        return _extract_from_docx(file_content)
    elif ext == '.pdf':
        return _extract_from_pdf(file_content)
    elif ext == '.pptx':
        return _extract_from_pptx(file_content)
    elif ext == '.xlsx':
        return _extract_from_xlsx(file_content)
    elif ext in ['.txt', '.md', '.markdown', '.json', '.yaml', '.yml']:
        return file_content.decode('utf-8', errors='ignore')
    else:
        return ""

def _extract_from_docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        log_print(f"Error reading DOCX: {e}")
        return ""

def _extract_from_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        full_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                full_text.append(text)
        return '\n'.join(full_text)
    except Exception as e:
        log_print(f"Error reading PDF: {e}")
        return ""

def _extract_from_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        full_text = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                full_text.append("\n".join(slide_text))
        return '\n\n'.join(full_text)
    except Exception as e:
        log_print(f"Error reading PPTX: {e}")
        return ""

def _extract_from_xlsx(content: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), data_only=True)
        full_text = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            full_text.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                # Filter out None values and convert to string
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:
                    full_text.append(" | ".join(row_text))
        return '\n'.join(full_text)
    except Exception as e:
        log_print(f"Error reading XLSX: {e}")
        return ""
