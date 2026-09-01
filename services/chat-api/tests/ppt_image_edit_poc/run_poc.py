from __future__ import annotations

import argparse
import base64
import json
import mimetypes
import os
import random
import re
import sys
import time
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List

import httpx
from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
try:
    import cv2  # type: ignore
except Exception:
    cv2 = None
import numpy as np

from app.llm.providers.azure_gpt_image import AzureGptImageClient


ROOT = Path(__file__).resolve().parent
BACKEND_ROOT = ROOT.parents[1]
ENV_PATH = BACKEND_ROOT / ".env"
SCHEMA_PATH = ROOT / "slide_schema.json"
SLIDE_W = 1536
SLIDE_H = 864


class RunLogger:
    def __init__(self, path: Path) -> None:
        self.path = path
        self.events: List[Dict[str, Any]] = []

    def log(self, event: str, **fields: Any) -> None:
        record = {
            "ts": datetime.now().isoformat(timespec="seconds"),
            "event": event,
            **sanitize_for_log(fields),
        }
        self.events.append(record)
        line = json.dumps(record, ensure_ascii=False)
        with self.path.open("a", encoding="utf-8") as f:
            f.write(line + "\n")
        print(f"[poc] {event} {json.dumps(sanitize_for_console(fields), ensure_ascii=False)}", flush=True)


def sanitize_for_log(value: Any) -> Any:
    if isinstance(value, dict):
        return {str(k): sanitize_for_log(v) for k, v in value.items() if str(k).lower() not in {"api-key", "authorization"}}
    if isinstance(value, list):
        return [sanitize_for_log(v) for v in value]
    if isinstance(value, tuple):
        return [sanitize_for_log(v) for v in value]
    if isinstance(value, str):
        if value.startswith("data:image/") and "base64," in value:
            return value.split("base64,", 1)[0] + "base64,<redacted>"
        if len(value) > 1600:
            return value[:1600] + f"...<truncated {len(value)} chars>"
    return value


def sanitize_for_console(value: Any) -> Any:
    compact = sanitize_for_log(value)
    if isinstance(compact, dict):
        return {k: sanitize_for_console(v) for k, v in compact.items()}
    if isinstance(compact, list):
        return [sanitize_for_console(v) for v in compact]
    if isinstance(compact, str) and len(compact) > 300:
        return compact[:300] + "...<truncated>"
    return compact


def env_summary(env: Dict[str, str]) -> Dict[str, Any]:
    return {
        "AZURE_IMAGE_API_STYLE": env.get("AZURE_IMAGE_API_STYLE", ""),
        "AZURE_IMAGE_GENERATION_API_STYLE": env.get("AZURE_IMAGE_GENERATION_API_STYLE", ""),
        "AZURE_IMAGE_EDIT_API_STYLE": env.get("AZURE_IMAGE_EDIT_API_STYLE", ""),
        "AZURE_IMAGE_OPENAI_ENDPOINT": env.get("AZURE_IMAGE_OPENAI_ENDPOINT", ""),
        "AZURE_IMAGE_OPENAI_API_VERSION": env.get("AZURE_IMAGE_OPENAI_API_VERSION", ""),
        "AZURE_IMAGE_EDIT_API_VERSION": env.get("AZURE_IMAGE_EDIT_API_VERSION", ""),
        "AZURE_IMAGE_DEPLOYMENT_NAME": env.get("AZURE_IMAGE_DEPLOYMENT_NAME", ""),
        "AZURE_IMAGE_SIZE": env.get("AZURE_IMAGE_SIZE", ""),
        "AZURE_IMAGE_QUALITY": env.get("AZURE_IMAGE_QUALITY", ""),
        "AZURE_IMAGE_MAX_RETRIES": env.get("AZURE_IMAGE_MAX_RETRIES", ""),
        "AZURE_IMAGE_RETRY_BASE_SECONDS": env.get("AZURE_IMAGE_RETRY_BASE_SECONDS", ""),
        "AZURE_IMAGE_RETRY_MAX_SECONDS": env.get("AZURE_IMAGE_RETRY_MAX_SECONDS", ""),
        "AZURE_IMAGE_HTTP_READ_TIMEOUT": env.get("AZURE_IMAGE_HTTP_READ_TIMEOUT", ""),
        "AZURE_IMAGE_HTTP_KEEPALIVE": env.get("AZURE_IMAGE_HTTP_KEEPALIVE", ""),
        "AZURE_IMAGE_V1_INCLUDE_API_VERSION": env.get("AZURE_IMAGE_V1_INCLUDE_API_VERSION", ""),
        "AZURE_IMAGE_OPENAI_API_KEY_present": bool(env.get("AZURE_IMAGE_OPENAI_API_KEY", "").strip()),
        "VISION_BASE_URL": env.get("VISION_BASE_URL", ""),
        "VISION_MODEL": env.get("VISION_MODEL", ""),
        "DASHSCOPE_API_KEY_present": bool(env.get("DASHSCOPE_API_KEY", "").strip()),
        "QWEN_API_KEY_present": bool(env.get("QWEN_API_KEY", "").strip()),
    }


def load_env(path: Path) -> Dict[str, str]:
    values = dict(os.environ)
    if not path.exists():
        return values
    for raw in path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip().strip('"').strip("'")
        if key and key not in values:
            values[key] = value
    return values


def read_schema() -> Dict[str, Any]:
    return json.loads(SCHEMA_PATH.read_text(encoding="utf-8"))


def slide_prompt(slide: Dict[str, Any]) -> str:
    text_lines = [f"- {slot['text']}" for slot in slide.get("text_slots", [])]
    chart_lines = []
    for chart in slide.get("chart_slots", []):
        labels = ", ".join(str(x) for x in chart.get("labels", []))
        values = ", ".join(str(x) for x in chart.get("values", []))
        chart_lines.append(f"- {chart.get('title')}: labels={labels}; values={values}")
    return (
        "Create one complete 16:9 PowerPoint slide visual draft.\n"
        "The output should look like a premium executive presentation slide, not a poster.\n"
        "Include the exact text below in the visual draft because this is draft A.\n"
        "Use refined typography, icons, background, lighting, visual hierarchy, and composition.\n"
        "Do not add extra unrelated text.\n\n"
        f"Visual style: {slide.get('visual_prompt', '')}\n\n"
        "Required visible text:\n"
        + "\n".join(text_lines)
        + "\n\nRequired chart/table content if applicable:\n"
        + ("\n".join(chart_lines) if chart_lines else "- none")
    )


def safe_name(value: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9_-]+", "_", str(value or "").strip()).strip("_")
    return text or "slide"


def parse_csv_tokens(raw: str) -> List[str]:
    if not str(raw or "").strip():
        return []
    return [part.strip() for part in str(raw).split(",") if part.strip()]


def ensure_run_dir() -> Path:
    run_dir = ROOT / "runs" / datetime.now().strftime("%Y%m%d_%H%M%S")
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


def latest_existing_slide_artifact(pattern: str) -> Path:
    candidates = sorted((ROOT / "runs").glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True)
    if not candidates:
        raise FileNotFoundError(f"no existing artifact found for pattern runs/{pattern}")
    return candidates[0]


def decode_image_response(data: Dict[str, Any], out_path: Path, env: Dict[str, str], log: RunLogger | None = None) -> None:
    items = data.get("data") or []
    if not items:
        raise RuntimeError(f"image response has no data: {str(data)[:800]}")
    first = items[0]
    if first.get("b64_json"):
        image_bytes = base64.b64decode(first["b64_json"])
        out_path.write_bytes(image_bytes)
        if log:
            log.log("image_decode_b64_done", out_path=str(out_path), bytes=len(image_bytes))
        return
    url = str(first.get("url") or "").strip()
    if not url:
        raise RuntimeError(f"image response has no b64_json/url: {str(data)[:800]}")
    started = time.monotonic()
    if log:
        log.log("image_download_start", url=url[:260], out_path=str(out_path))
    with httpx.Client(timeout=90, follow_redirects=True) as client:
        resp = client.get(url)
        if log:
            log.log(
                "image_download_response",
                status_code=resp.status_code,
                elapsed_ms=int((time.monotonic() - started) * 1000),
                content_type=resp.headers.get("content-type", ""),
                bytes=len(resp.content or b""),
            )
        resp.raise_for_status()
        out_path.write_bytes(resp.content)


def image_size_str(path: Path) -> str:
    with Image.open(path) as img:
        w, h = img.size
    return f"{w}x{h}"


def ensure_official_edit_mask(mask_path: Path, image_path: Path, out_path: Path | None = None) -> Path:
    """
    Azure GPT Image edit masks require a PNG mask with an alpha channel and the same
    dimensions as the input image. Fully transparent pixels (alpha=0) define the edit
    area, so we invert our internal white-text mask for the API payload.
    """
    target = out_path or mask_path
    with Image.open(image_path) as img:
        image_size = img.size
    mask_l = Image.open(mask_path).convert("L")
    if mask_l.size != image_size:
        mask_l = mask_l.resize(image_size, Image.Resampling.LANCZOS)
    alpha = Image.eval(mask_l, lambda px: 0 if px > 0 else 255)
    mask_rgba = Image.new("RGBA", image_size, (255, 255, 255, 255))
    mask_rgba.putalpha(alpha)
    mask_rgba.save(target, format="PNG")
    return target


def mask_alpha_stats(mask_path: Path) -> Dict[str, Any]:
    mask = Image.open(mask_path).convert("RGBA")
    alpha = mask.getchannel("A")
    hist = alpha.histogram()
    total = mask.size[0] * mask.size[1]
    transparent = hist[0]
    nontransparent = total - transparent
    return {
        "mode": mask.mode,
        "size": f"{mask.size[0]}x{mask.size[1]}",
        "transparent_pixels": transparent,
        "transparent_ratio": round(transparent / total, 4) if total else 0,
        "nontransparent_pixels": nontransparent,
        "nontransparent_ratio": round(nontransparent / total, 4) if total else 0,
        "alpha_nonzero_bbox": alpha.getbbox(),
    }


def stabilize_edited_background(original_path: Path, edited_path: Path, mask_path: Path, log: RunLogger | None = None) -> None:
    """
    Keep non-masked regions identical to original A image.
    This prevents edit-model artifacts (missing/cropped/black regions) outside text-mask areas.
    """
    original_rgba = Image.open(original_path).convert("RGBA")
    edited_rgba = Image.open(edited_path).convert("RGBA")
    if edited_rgba.size != original_rgba.size:
        edited_rgba = edited_rgba.resize(original_rgba.size, Image.Resampling.LANCZOS)
    raw_mask = Image.open(mask_path)
    if "A" in raw_mask.getbands():
        # For RGBA masks, edit region is encoded in alpha channel.
        mask_l = raw_mask.convert("RGBA").getchannel("A")
    else:
        mask_l = raw_mask.convert("L")
    if mask_l.size != original_rgba.size:
        mask_l = mask_l.resize(original_rgba.size, Image.Resampling.LANCZOS)
    # Feather edges to reduce seams at text boundaries.
    region_mask = mask_l.filter(ImageFilter.GaussianBlur(radius=2))
    candidate_rgba = Image.alpha_composite(original_rgba, edited_rgba)
    merged = Image.composite(candidate_rgba.convert("RGB"), original_rgba.convert("RGB"), region_mask)
    merged.save(edited_path)
    if log:
        log.log(
            "image_edit_stabilized",
            original=str(original_path),
            edited=str(edited_path),
            mask=str(mask_path),
            size=f"{original_rgba.size[0]}x{original_rgba.size[1]}",
        )


def normalize_slide_image(path: Path) -> None:
    img = Image.open(path).convert("RGB")
    if img.size == (SLIDE_W, SLIDE_H):
        img.save(path)
        return
    src_w, src_h = img.size
    target_ratio = SLIDE_W / SLIDE_H
    src_ratio = src_w / src_h
    if src_ratio > target_ratio:
        new_w = int(src_h * target_ratio)
        left = max(0, (src_w - new_w) // 2)
        img = img.crop((left, 0, left + new_w, src_h))
    else:
        new_h = int(src_w / target_ratio)
        top = max(0, (src_h - new_h) // 2)
        img = img.crop((0, top, src_w, top + new_h))
    img = img.resize((SLIDE_W, SLIDE_H), Image.Resampling.LANCZOS)
    img.save(path)


def image_generation_endpoint(env: Dict[str, str], path: str, *, operation: str = "") -> str:
    endpoint = env.get("AZURE_IMAGE_OPENAI_ENDPOINT", "").rstrip("/")
    if not endpoint:
        raise RuntimeError("AZURE_IMAGE_OPENAI_ENDPOINT is not configured")
    deployment = env.get("AZURE_IMAGE_DEPLOYMENT_NAME", "gpt-image-2").strip() or "gpt-image-2"
    default_style = env.get("AZURE_IMAGE_API_STYLE", "").strip().lower()
    include_version_in_v1 = str(env.get("AZURE_IMAGE_V1_INCLUDE_API_VERSION", "")).strip().lower() in {"1", "true", "yes"}
    if operation == "edit":
        # Historical stable path in this PoC: /openai/v1/images/edits (without api-version query).
        style = env.get("AZURE_IMAGE_EDIT_API_STYLE", "").strip().lower() or default_style or "v1"
        api_version = env.get("AZURE_IMAGE_EDIT_API_VERSION", "").strip() or env.get("AZURE_IMAGE_OPENAI_API_VERSION", "").strip() or "2025-04-01-preview"
        if style == "v1":
            if include_version_in_v1:
                return f"{endpoint}/openai/v1/images/{path}?api-version={api_version}"
            return f"{endpoint}/openai/v1/images/{path}"
        return f"{endpoint}/openai/deployments/{deployment}/images/{path}?api-version={api_version}"
    # Historical stable path in this PoC: /openai/v1/images/generations (without api-version query).
    style = env.get("AZURE_IMAGE_GENERATION_API_STYLE", "").strip().lower() or default_style or "v1"
    api_version = env.get("AZURE_IMAGE_OPENAI_API_VERSION", "").strip() or "2024-02-01"
    if style == "v1":
        if include_version_in_v1:
            return f"{endpoint}/openai/v1/images/{path}?api-version={api_version}"
        return f"{endpoint}/openai/v1/images/{path}"
    return f"{endpoint}/openai/deployments/{deployment}/images/{path}?api-version={api_version}"


def env_int(env: Dict[str, str], key: str, default: int) -> int:
    raw = str(env.get(key, "")).strip()
    if not raw:
        return default
    try:
        return int(raw)
    except Exception:
        return default


def env_float(env: Dict[str, str], key: str, default: float) -> float:
    raw = str(env.get(key, "")).strip()
    if not raw:
        return default
    try:
        return float(raw)
    except Exception:
        return default


def image_http_timeout(env: Dict[str, str]) -> httpx.Timeout:
    connect = env_float(env, "AZURE_IMAGE_HTTP_CONNECT_TIMEOUT", 30.0)
    read = env_float(env, "AZURE_IMAGE_HTTP_READ_TIMEOUT", 600.0)
    write = env_float(env, "AZURE_IMAGE_HTTP_WRITE_TIMEOUT", 600.0)
    pool = env_float(env, "AZURE_IMAGE_HTTP_POOL_TIMEOUT", 30.0)
    return httpx.Timeout(connect=connect, read=read, write=write, pool=pool)


def image_http_limits(env: Dict[str, str]) -> httpx.Limits:
    keepalive = max(0, env_int(env, "AZURE_IMAGE_HTTP_KEEPALIVE", 0))
    connections = max(1, env_int(env, "AZURE_IMAGE_HTTP_MAX_CONNECTIONS", 2))
    return httpx.Limits(max_keepalive_connections=keepalive, max_connections=connections)


def validate_image_size(size: str, deployment: str) -> None:
    raw = str(size or "").strip().lower()
    if raw == "auto":
        return
    if not re.fullmatch(r"\d+x\d+", raw):
        raise ValueError(f"invalid image size '{size}', expected WIDTHxHEIGHT or auto")
    w, h = [int(x) for x in raw.split("x", 1)]
    model = str(deployment or "").lower()
    if "gpt-image-2" in model:
        pixels = w * h
        if w % 16 != 0 or h % 16 != 0:
            raise ValueError(f"gpt-image-2 size '{size}' must have width and height divisible by 16")
        if pixels < 65_536 or pixels > 4_194_304:
            raise ValueError(f"gpt-image-2 size '{size}' must be between 65,536 and 4,194,304 pixels")
        ratio = w / h
        if ratio < 0.5 or ratio > 2.0:
            raise ValueError(f"gpt-image-2 size '{size}' aspect ratio must be between 1:2 and 2:1")
        return
    allowed = {"1024x1024", "1024x1536", "1536x1024"}
    if raw not in allowed:
        raise ValueError(f"image size '{size}' is not in the common Azure GPT Image size set {sorted(allowed)} for deployment '{deployment}'")


def parse_retry_after_seconds(resp: httpx.Response) -> float | None:
    for key in ("retry-after-ms", "x-ms-retry-after-ms"):
        raw_ms = str(resp.headers.get(key, "")).strip()
        if raw_ms:
            try:
                return max(0.0, float(raw_ms) / 1000.0)
            except Exception:
                pass
    raw = str(resp.headers.get("retry-after", "")).strip()
    if not raw:
        return None
    try:
        return max(0.0, float(raw))
    except Exception:
        return None


def should_retry_status(code: int) -> bool:
    return code in {408, 409, 425, 429, 500, 502, 503, 504}


def post_with_retry(
    endpoint: str,
    *,
    headers: Dict[str, str],
    env: Dict[str, str],
    log: RunLogger | None,
    request_tag: str,
    json_payload: Dict[str, Any] | None = None,
    form_data: Dict[str, Any] | None = None,
    files: Dict[str, Any] | None = None,
) -> httpx.Response:
    max_retries = max(0, env_int(env, "AZURE_IMAGE_MAX_RETRIES", 5))
    max_attempts = max_retries + 1
    base_delay = max(0.2, env_float(env, "AZURE_IMAGE_RETRY_BASE_SECONDS", 1.5))
    cap_delay = max(base_delay, env_float(env, "AZURE_IMAGE_RETRY_MAX_SECONDS", 30.0))
    last_exc: Exception | None = None

    for attempt in range(1, max_attempts + 1):
        started = time.monotonic()
        try:
            with httpx.Client(timeout=image_http_timeout(env), limits=image_http_limits(env), follow_redirects=True, http2=False) as client:
                resp = client.post(endpoint, headers=headers, json=json_payload, data=form_data, files=files)
        except (httpx.ConnectError, httpx.ReadError, httpx.RemoteProtocolError, httpx.TimeoutException) as exc:
            last_exc = exc
            retryable = attempt < max_attempts
            if log:
                log.log(
                    f"{request_tag}_transport_error",
                    endpoint=endpoint,
                    attempt=attempt,
                    max_attempts=max_attempts,
                    elapsed_ms=int((time.monotonic() - started) * 1000),
                    exception_type=type(exc).__name__,
                    error=str(exc),
                    retryable=retryable,
                )
            if not retryable:
                raise
            delay = min(cap_delay, base_delay * (2 ** (attempt - 1)))
            delay = min(cap_delay, delay + random.uniform(0, delay * 0.25))
            if log:
                log.log(f"{request_tag}_retry_sleep", endpoint=endpoint, attempt=attempt, sleep_seconds=round(delay, 3), reason="transport_error")
            time.sleep(delay)
            continue

        status = resp.status_code
        if status < 400:
            return resp
        if not should_retry_status(status) or attempt >= max_attempts:
            return resp
        retry_after = parse_retry_after_seconds(resp)
        delay = retry_after if retry_after is not None else min(cap_delay, base_delay * (2 ** (attempt - 1)))
        delay = min(cap_delay, delay + random.uniform(0, max(0.05, delay * 0.25)))
        if log:
            log.log(
                f"{request_tag}_retryable_status",
                endpoint=endpoint,
                attempt=attempt,
                max_attempts=max_attempts,
                status_code=status,
                retry_after_seconds=retry_after,
                sleep_seconds=round(delay, 3),
                response_preview=(resp.text or "")[:500],
            )
        time.sleep(delay)

    if last_exc:
        raise last_exc
    raise RuntimeError(f"{request_tag} failed without response")


def generate_image_api(prompt: str, out_path: Path, env: Dict[str, str], log: RunLogger | None = None) -> Dict[str, Any]:
    client = AzureGptImageClient(env=env)
    result = client.generate_image(
        prompt,
        log_hook=(lambda event, payload: log.log(event, **payload)) if log else None,
    )
    out_path.write_bytes(result.image_bytes)
    return {"ok": True, "endpoint": result.endpoint, "response": result.response}


def edit_image_api(image_path: Path, mask_path: Path | None, prompt: str, out_path: Path, env: Dict[str, str], log: RunLogger | None = None) -> Dict[str, Any]:
    api_key = env.get("AZURE_IMAGE_OPENAI_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("missing AZURE_IMAGE_OPENAI_API_KEY")
    endpoint = image_generation_endpoint(env, "edits", operation="edit")
    deployment = env.get("AZURE_IMAGE_DEPLOYMENT_NAME", "gpt-image-2").strip() or "gpt-image-2"
    endpoint_style = env.get("AZURE_IMAGE_EDIT_API_STYLE", "").strip().lower() or env.get("AZURE_IMAGE_API_STYLE", "").strip().lower() or "deployments"
    default_auth_style = "api-key" if endpoint_style == "v1" else "bearer"
    auth_style = str(env.get("AZURE_IMAGE_EDIT_AUTH_STYLE", default_auth_style)).strip().lower()
    headers = {"Connection": "close"}
    if auth_style == "api-key":
        headers["api-key"] = api_key
    else:
        headers["Authorization"] = f"Bearer {api_key}"
    edit_size = image_size_str(image_path)
    validate_image_size(edit_size, deployment)
    official_mask_path: Path | None = None
    source_mask_stats: Dict[str, Any] | None = None
    api_mask_stats: Dict[str, Any] | None = None
    mask_note: Dict[str, Any]
    if mask_path is not None:
        # Azure edit commonly treats transparent alpha as editable area.
        # Our internal mask marks text with non-zero alpha, so default to inversion for API payload.
        convert_mask_default = True
        convert_mask = str(env.get("AZURE_IMAGE_EDIT_CONVERT_MASK_ALPHA", str(convert_mask_default))).strip().lower() in {"1", "true", "yes"}
        source_mask_stats = mask_alpha_stats(mask_path)
        if convert_mask:
            official_mask_path = ensure_official_edit_mask(mask_path, image_path, mask_path.with_name(f"{mask_path.stem}_api.png"))
            mask_note = {"converted_alpha_mask": True, "transparent_pixels_are_edit_area": True}
        else:
            official_mask_path = mask_path
            mask_note = {"converted_alpha_mask": False, "mask_sent_as_provided": True}
        api_mask_stats = mask_alpha_stats(official_mask_path)
    else:
        mask_note = {"mask_omitted": True, "edit_strategy": "prompt_only"}
    image_bytes = image_path.read_bytes()
    mask_bytes = official_mask_path.read_bytes() if official_mask_path is not None else b""
    # Historical stable behavior in this PoC used "image" form field.
    default_image_field = "image"
    image_field = env.get("AZURE_IMAGE_EDIT_IMAGE_FIELD", default_image_field).strip() or default_image_field
    files = {
        image_field: (image_path.name, image_bytes, mimetypes.guess_type(str(image_path))[0] or "image/png"),
    }
    if official_mask_path is not None:
        files["mask"] = (official_mask_path.name, mask_bytes, "image/png")
    data = {"prompt": prompt}
    include_model_default = endpoint_style == "v1"
    include_model = str(env.get("AZURE_IMAGE_EDIT_INCLUDE_MODEL", str(include_model_default))).strip().lower() in {"1", "true", "yes"}
    if include_model:
        data["model"] = deployment
    include_size_default = endpoint_style == "v1"
    if str(env.get("AZURE_IMAGE_EDIT_INCLUDE_SIZE", str(include_size_default))).strip().lower() in {"1", "true", "yes"}:
        data["size"] = edit_size
    include_quality_default = endpoint_style == "v1"
    if str(env.get("AZURE_IMAGE_EDIT_INCLUDE_QUALITY", str(include_quality_default))).strip().lower() in {"1", "true", "yes"}:
        data["quality"] = env.get("AZURE_IMAGE_QUALITY", "low")
    include_output_format_default = endpoint_style == "v1"
    if str(env.get("AZURE_IMAGE_EDIT_INCLUDE_OUTPUT_FORMAT", str(include_output_format_default))).strip().lower() in {"1", "true", "yes"}:
        data["output_format"] = "png"
    include_n_default = endpoint_style == "v1"
    if str(env.get("AZURE_IMAGE_EDIT_INCLUDE_N", str(include_n_default))).strip().lower() in {"1", "true", "yes"}:
        data["n"] = "1"
    started = time.monotonic()
    if log:
        log.log(
            "image_edit_request",
            endpoint=endpoint,
            payload={
                "fields": sorted(data.keys()),
                "size": data.get("size", "<omitted>"),
                "model": data.get("model", "<omitted>"),
                "quality": data.get("quality", "<omitted>"),
                "output_format": data.get("output_format", "<omitted>"),
                "n": data.get("n", "<omitted>"),
                "prompt_chars": len(prompt),
                "prompt_preview": prompt[:700],
            },
            auth_style=auth_style,
            image={"path": str(image_path), "bytes": len(image_bytes), "content_type": files[image_field][2], "field": image_field},
            mask=(
                {
                    "path": str(official_mask_path),
                    "bytes": len(mask_bytes),
                    "content_type": "image/png",
                    "same_size_as_image": image_size_str(official_mask_path) == edit_size,
                    "source_mask_alpha": source_mask_stats,
                    "api_mask_alpha": api_mask_stats,
                    **mask_note,
                }
                if official_mask_path is not None
                else mask_note
            ),
        )
    if log:
        log.log(
            "image_edit_http_config",
            timeout={
                "connect": env_float(env, "AZURE_IMAGE_HTTP_CONNECT_TIMEOUT", 30.0),
                "read": env_float(env, "AZURE_IMAGE_HTTP_READ_TIMEOUT", 600.0),
                "write": env_float(env, "AZURE_IMAGE_HTTP_WRITE_TIMEOUT", 600.0),
                "pool": env_float(env, "AZURE_IMAGE_HTTP_POOL_TIMEOUT", 30.0),
            },
            keepalive=env_int(env, "AZURE_IMAGE_HTTP_KEEPALIVE", 0),
        )
    try:
        resp = post_with_retry(
            endpoint,
            headers=headers,
            env=env,
            log=log,
            request_tag="image_edit",
            form_data=data,
            files=files,
        )
    except Exception as exc:
        if log:
            log.log(
                "image_edit_exception",
                endpoint=endpoint,
                elapsed_ms=int((time.monotonic() - started) * 1000),
                exception_type=type(exc).__name__,
                error=str(exc),
            )
        raise
    if log:
        log.log(
            "image_edit_response",
            endpoint=endpoint,
            status_code=resp.status_code,
            elapsed_ms=int((time.monotonic() - started) * 1000),
            content_type=resp.headers.get("content-type", ""),
            request_id=resp.headers.get("x-ms-request-id", "") or resp.headers.get("apim-request-id", ""),
            response_preview=(resp.text or "")[:1200],
        )
    if resp.status_code >= 400:
        raise RuntimeError(f"edit http_{resp.status_code}: {resp.text[:1200]}")
    payload = resp.json()
    decode_image_response(payload, out_path, env, log)
    if mask_path is not None:
        stabilize_edited_background(image_path, out_path, mask_path, log)
    return {"ok": True, "endpoint": endpoint, "response": compact_response(payload)}


def compact_response(data: Dict[str, Any]) -> Dict[str, Any]:
    out = dict(data)
    items = []
    for item in out.get("data") or []:
        compact = dict(item)
        if compact.get("b64_json"):
            compact["b64_json"] = f"<base64:{len(compact['b64_json'])} chars>"
        if compact.get("url"):
            compact["url"] = str(compact["url"])[:240]
        items.append(compact)
    if items:
        out["data"] = items
    return out


def contains_cjk(text: str) -> bool:
    for ch in str(text or ""):
        code = ord(ch)
        if 0x3400 <= code <= 0x9FFF or 0xF900 <= code <= 0xFAFF:
            return True
    return False


def normalize_font_family_hint(value: Any, default: str = "sans") -> str:
    token = str(value or "").strip().lower()
    if not token:
        return default
    if any(k in token for k in ["mono", "code", "menlo", "consola", "courier"]):
        return "mono"
    if any(k in token for k in ["serif", "song", "times", "georgia"]):
        return "serif"
    if any(k in token for k in ["sans", "inter", "arial", "helvetica", "pingfang", "yahei", "noto"]):
        return "sans"
    return default


def parse_font_weight(value: Any, default_token: str = "regular") -> tuple[str, int]:
    raw = str(value or "").strip().lower()
    if not raw:
        return ("bold", 700) if default_token == "bold" else ("regular", 400)
    if re.fullmatch(r"\d{3}", raw):
        w = max(100, min(900, int(raw)))
        return ("bold" if w >= 600 else "regular", w)
    if raw in {"bold", "bolder", "semibold", "demibold"}:
        return ("bold", 700 if raw == "bold" else 600)
    if raw in {"regular", "normal", "book", "medium", "light"}:
        return ("regular", 400 if raw in {"regular", "normal", "book"} else (500 if raw == "medium" else 300))
    return ("bold", 700) if default_token == "bold" else ("regular", 400)


def font_targets(text: str, family_hint: str, weight_value: int, bold: bool) -> Dict[str, Any]:
    is_cjk = contains_cjk(text)
    fam = normalize_font_family_hint(family_hint, "sans")
    heavy = bool(bold or weight_value >= 600)
    if fam == "serif":
        if is_cjk:
            return {
                "pil_paths": ["/System/Library/Fonts/Supplemental/Songti.ttc", "/System/Library/Fonts/STHeiti Light.ttc", "/System/Library/Fonts/PingFang.ttc"],
                "ppt_name": "Songti SC",
            }
        return {
            "pil_paths": ["/System/Library/Fonts/Supplemental/Times New Roman Bold.ttf" if heavy else "/System/Library/Fonts/Supplemental/Times New Roman.ttf", "/System/Library/Fonts/Supplemental/Times.ttf"],
            "ppt_name": "Times New Roman",
        }
    if fam == "mono":
        if is_cjk:
            return {
                "pil_paths": ["/System/Library/Fonts/PingFang.ttc", "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"],
                "ppt_name": "PingFang SC",
            }
        return {
            "pil_paths": ["/System/Library/Fonts/Menlo.ttc", "/System/Library/Fonts/Supplemental/Courier New.ttf"],
            "ppt_name": "Menlo",
        }
    # sans default
    if is_cjk:
        return {
            "pil_paths": ["/System/Library/Fonts/PingFang.ttc", "/System/Library/Fonts/Supplemental/Arial Unicode.ttf", "/System/Library/Fonts/Supplemental/Arial.ttf"],
            "ppt_name": "PingFang SC",
        }
    return {
        "pil_paths": ["/System/Library/Fonts/Supplemental/Arial Bold.ttf" if heavy else "/System/Library/Fonts/Supplemental/Arial.ttf", "/Library/Fonts/Arial Unicode.ttf"],
        "ppt_name": "Arial",
    }


def font(size: int, bold: bool = False, *, text: str = "", family_hint: str = "sans", weight_value: int = 400) -> ImageFont.ImageFont:
    targets = font_targets(text, family_hint, weight_value, bold)
    candidates = list(targets.get("pil_paths") or []) + ["/Library/Fonts/Arial Unicode.ttf"]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            try:
                return ImageFont.truetype(candidate, size=size)
            except Exception:
                pass
    return ImageFont.load_default()


def wrap_text_to_width(
    text: str,
    max_width: int,
    fnt: ImageFont.ImageFont,
    *,
    allow_wrap: bool = True,
    prefer_balanced_two_lines: bool = False,
) -> str:
    content = str(text or "").strip()
    if not content:
        return ""
    if not allow_wrap:
        return content.replace("\n", " ")
    lines: List[str] = []
    draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    for paragraph in content.split("\n"):
        paragraph = paragraph.strip()
        if not paragraph:
            lines.append("")
            continue
        if prefer_balanced_two_lines and " " in paragraph:
            words = [w for w in paragraph.split(" ") if w]
            if len(words) >= 4:
                best: tuple[float, str, str] | None = None
                for i in range(1, len(words)):
                    left = " ".join(words[:i])
                    right = " ".join(words[i:])
                    left_bbox = draw.textbbox((0, 0), left, font=fnt)
                    right_bbox = draw.textbbox((0, 0), right, font=fnt)
                    left_w = left_bbox[2] - left_bbox[0]
                    right_w = right_bbox[2] - right_bbox[0]
                    if left_w <= max_width and right_w <= max_width:
                        score = abs(left_w - right_w) + abs(len(words[:i]) - len(words[i:])) * 6
                        if best is None or score < best[0]:
                            best = (score, left, right)
                if best is not None:
                    lines.append(best[1])
                    lines.append(best[2])
                    continue
        # Prefer word wrapping for space-delimited languages.
        if " " in paragraph:
            words = [w for w in paragraph.split(" ") if w]
            if not words:
                lines.append("")
                continue
            cur = words[0]
            for word in words[1:]:
                candidate = f"{cur} {word}"
                bbox = draw.textbbox((0, 0), candidate, font=fnt)
                if (bbox[2] - bbox[0]) <= max_width:
                    cur = candidate
                else:
                    lines.append(cur)
                    cur = word
            lines.append(cur)
            continue
        # Fallback char wrapping for no-space text (e.g., CJK).
        cur = ""
        for ch in paragraph:
            candidate = cur + ch
            bbox = draw.textbbox((0, 0), candidate, font=fnt)
            if cur and (bbox[2] - bbox[0]) > max_width:
                lines.append(cur)
                cur = ch
            else:
                cur = candidate
        if cur:
            lines.append(cur)
    return "\n".join(lines)


def fit_wrapped_text(
    text: str,
    box_w: int,
    box_h: int,
    requested: int,
    bold: bool = False,
    *,
    allow_wrap: bool = True,
    prefer_balanced_two_lines: bool = False,
    family_hint: str = "sans",
    weight_value: int = 400,
) -> tuple[str, int]:
    size = max(10, int(requested))
    draw = ImageDraw.Draw(Image.new("RGB", (1, 1)))
    best_text = str(text or "")
    best_size = size
    while size > 10:
        fnt = font(size, bold, text=str(text or ""), family_hint=family_hint, weight_value=weight_value)
        wrapped = wrap_text_to_width(
            text,
            max(1, box_w),
            fnt,
            allow_wrap=allow_wrap,
            prefer_balanced_two_lines=prefer_balanced_two_lines,
        )
        bbox = draw.multiline_textbbox((0, 0), wrapped, font=fnt, spacing=6)
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]
        if text_w <= max(1, box_w) and text_h <= max(1, box_h):
            return wrapped, size
        best_text = wrapped
        best_size = size
        size -= 1
    return best_text, best_size


def draw_slot_text(draw: ImageDraw.ImageDraw, slot: Dict[str, Any], text: str, scale: float = 1.5, *, force_no_wrap: bool = False) -> None:
    x, y, w, h = rect_px(slot)
    weight_token, weight_value = parse_font_weight(slot.get("font_weight"), "regular")
    if slot.get("font_weight_value") is not None:
        _, weight_value = parse_font_weight(slot.get("font_weight_value"), weight_token)
    bold = weight_token == "bold" or weight_value >= 600
    family_hint = normalize_font_family_hint(slot.get("font_family_hint"), "sans")
    requested = int(float(slot.get("font_size", 18)) * scale)
    role = str(slot.get("role") or "").strip().lower()
    allow_wrap = (role not in {"subtitle", "caption", "footer"}) and (not force_no_wrap)
    wrapped, fitted = fit_wrapped_text(
        text,
        w,
        h,
        requested,
        bold,
        allow_wrap=allow_wrap,
        prefer_balanced_two_lines=(role == "title"),
        family_hint=family_hint,
        weight_value=weight_value,
    )
    draw.multiline_text(
        (x, y),
        wrapped,
        fill=hex_to_rgba(str(slot.get("color", "#ffffff")), 255),
        font=font(fitted, bold, text=str(text or ""), family_hint=family_hint, weight_value=weight_value),
        spacing=6,
        align=str(slot.get("align", "left")),
    )


def draw_line_spans(draw: ImageDraw.ImageDraw, line_slot: Dict[str, Any], spans: List[Dict[str, Any]], *, scale: float = 1.0) -> None:
    x, y, w, _ = rect_px(line_slot)
    align = str(line_slot.get("align", "left")).strip().lower()
    ordered = sorted(
        [s for s in spans if s.get("text") is not None and str(s.get("text")) != ""],
        key=lambda s: (
            int(s.get("order", 0)) if str(s.get("order", "")).isdigit() else 0,
            float((s.get("bbox") or {}).get("x", 0)),
        ),
    )
    if not ordered:
        draw_slot_text(draw, line_slot, str(line_slot.get("text") or ""), scale=scale, force_no_wrap=True)
        return
    segments: List[Dict[str, Any]] = []
    total_w = 0
    for span in ordered:
        text = str(span.get("text") or "")
        token, weight_value = parse_font_weight(span.get("font_weight"), "regular")
        if span.get("font_weight_value") is not None:
            _, weight_value = parse_font_weight(span.get("font_weight_value"), token)
        bold = token == "bold" or weight_value >= 600
        family_hint = normalize_font_family_hint(span.get("font_family_hint") or line_slot.get("font_family_hint"), "sans")
        requested = max(8, int(round(float(span.get("font_size") or line_slot.get("font_size") or 18) * scale)))
        fnt = font(requested, bold, text=text, family_hint=family_hint, weight_value=weight_value)
        bbox = draw.textbbox((0, 0), text, font=fnt)
        text_w = max(1, int(bbox[2] - bbox[0]))
        total_w += text_w
        segments.append({"text": text, "font": fnt, "width": text_w, "color": str(span.get("color") or line_slot.get("color") or "#ffffff")})
    start_x = x
    if align == "center":
        start_x = x + max(0, (w - total_w) // 2)
    elif align == "right":
        start_x = x + max(0, w - total_w)
    cursor = start_x
    for seg in segments:
        draw.text((cursor, y), seg["text"], fill=hex_to_rgba(seg["color"], 255), font=seg["font"])
        cursor += int(seg["width"])


def dry_run_image(slide: Dict[str, Any], out_path: Path) -> None:
    img = Image.new("RGB", (SLIDE_W, SLIDE_H), "#101827")
    draw = ImageDraw.Draw(img, "RGBA")
    for i in range(0, SLIDE_W, 48):
        alpha = int(25 * (1 - i / SLIDE_W))
        draw.line([(i, 0), (i - 260, SLIDE_H)], fill=(80, 160, 255, max(0, alpha)), width=2)
    draw.ellipse((980, -180, 1740, 580), fill=(34, 211, 238, 38))
    draw.ellipse((760, 470, 1320, 1040), fill=(134, 239, 172, 24))
    draw.rounded_rectangle((70, 70, SLIDE_W - 70, SLIDE_H - 70), radius=34, outline=(148, 163, 184, 55), width=2)
    for slot in slide.get("text_slots", []):
        x, y, w, h = rect_px(slot)
        draw.rounded_rectangle((x - 16, y - 10, x + w + 16, y + h + 10), radius=16, fill=(15, 23, 42, 85))
        draw_slot_text(draw, slot, str(slot.get("text", "")))
    for chart in slide.get("chart_slots", []):
        draw_chart_preview(draw, chart)
    img.save(out_path)


def rect_px(slot: Dict[str, Any]) -> tuple[int, int, int, int]:
    return (
        int(float(slot.get("x", 0)) * SLIDE_W),
        int(float(slot.get("y", 0)) * SLIDE_H),
        int(float(slot.get("w", 0)) * SLIDE_W),
        int(float(slot.get("h", 0)) * SLIDE_H),
    )


def hex_to_rgba(value: str, alpha: int) -> tuple[int, int, int, int]:
    raw = str(value or "#ffffff").strip().lstrip("#")
    if len(raw) != 6:
        return (255, 255, 255, alpha)
    return (int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16), alpha)


def draw_chart_preview(draw: ImageDraw.ImageDraw, chart: Dict[str, Any]) -> None:
    x, y, w, h = rect_px(chart)
    draw.rounded_rectangle((x, y, x + w, y + h), radius=26, fill=(2, 6, 23, 115), outline=(125, 211, 252, 80), width=2)
    labels = chart.get("labels", [])
    values = chart.get("values", [])
    if not labels or not values:
        return
    max_value = max(float(v) for v in values) or 1
    bar_w = max(24, int(w / (len(values) * 2.4)))
    base = y + h - 52
    for idx, value in enumerate(values):
        bx = x + 70 + idx * int(w / max(1, len(values)))
        bh = int((float(value) / max_value) * (h - 120))
        draw.rounded_rectangle((bx, base - bh, bx + bar_w, base), radius=10, fill=(34, 211, 238, 190))
        draw.text((bx - 18, base + 12), str(labels[idx])[:14], fill=(203, 213, 225, 230), font=font(18))


def default_layout(slide: Dict[str, Any]) -> Dict[str, Any]:
    layout = {
        "source": "schema_fallback",
        "text_slots": [
            {
                "id": slot["id"],
                "bbox": {"x": slot["x"], "y": slot["y"], "w": slot["w"], "h": slot["h"]},
                "align": slot.get("align", "left"),
                "font_size": slot.get("font_size", 20),
                "font_weight": slot.get("font_weight", "regular"),
                "font_weight_value": 700 if str(slot.get("font_weight", "")).strip().lower() == "bold" else 400,
                "font_family_hint": "sans",
                "color": slot.get("color", "#ffffff"),
                "role": slot.get("role", "body"),
            }
            for slot in slide.get("text_slots", [])
        ],
    }
    layout["text_lines"] = [
        {
            "line_id": f"{slot['id']}_l1",
            "order": 1,
            "slot_id": slot["id"],
            "text": str(slot.get("text", "")),
            "bbox": {"x": slot["x"], "y": slot["y"], "w": slot["w"], "h": slot["h"]},
            "align": slot.get("align", "left"),
            "font_size": slot.get("font_size", 20),
            "font_weight": slot.get("font_weight", "regular"),
            "font_weight_value": 700 if str(slot.get("font_weight", "")).strip().lower() == "bold" else 400,
            "font_family_hint": "sans",
            "color": slot.get("color", "#ffffff"),
            "confidence": 0.0,
        }
        for slot in slide.get("text_slots", [])
    ]
    layout["text_spans"] = [
        {
            "span_id": f"{slot['id']}_l1_s1",
            "line_id": f"{slot['id']}_l1",
            "slot_id": slot["id"],
            "order": 1,
            "text": str(slot.get("text", "")),
            "char_start": 0,
            "char_end": len(str(slot.get("text", ""))),
            "bbox": {"x": slot["x"], "y": slot["y"], "w": slot["w"], "h": slot["h"]},
            "font_size": slot.get("font_size", 20),
            "font_weight": slot.get("font_weight", "regular"),
            "font_weight_value": 700 if str(slot.get("font_weight", "")).strip().lower() == "bold" else 400,
            "font_family_hint": "sans",
            "color": slot.get("color", "#ffffff"),
            "confidence": 0.0,
        }
        for slot in slide.get("text_slots", [])
    ]
    return layout


def build_vlm_prompt(slide: Dict[str, Any]) -> str:
    slots = [
        {"id": s["id"], "role": s.get("role"), "schema_text": s.get("text")}
        for s in slide.get("text_slots", [])
    ]
    return (
        "You are locating editable text regions in a generated PowerPoint slide image.\n"
        "Do OCR and style extraction from the image.\n"
        "For each schema slot, infer text bounding box and style from the image.\n"
        "You must return one object for every schema slot id, with no omissions.\n"
        "font_size must be an integer in image pixels and must be measured per rendered line/span from visible glyph height, not copied or shared by slot.\n"
        "Use this definition: font_size is the typographic size whose uppercase glyph core height is closest to the observed glyph core height in the image.\n"
        "Estimate from glyph core only (ignore glow, shadow, blur, stroke, and outer effects).\n"
        "Different visual text heights must produce different font_size values, even within the same slot/title block.\n"
        "Do not force style consistency across lines: each line must have independently measured font_size.\n"
        "Also infer font_family_hint per item from this closed set: sans|serif|mono.\n"
        "Return font_weight as regular|bold and font_weight_value as numeric 100..900.\n"
        "Return line-level OCR results with one object per rendered line.\n"
        "Return span-level style results with one object per contiguous same-style segment in each line.\n"
        "If a line has mixed styles/colors, you must split into multiple spans and never collapse into dominant color.\n"
        "color must be core glyph fill color in strict #RRGGBB format (ignore glow/shadow).\n"
        "If uncertain, keep item and lower confidence instead of omitting.\n"
        "Do not copy style values from schema_text; infer from visual appearance in the image.\n"
        "Return strict JSON only with this shape: "
        "{\"text_slots\":[{\"id\":\"...\",\"bbox\":{\"x\":0,\"y\":0,\"w\":0,\"h\":0},"
        "\"align\":\"left|center|right\",\"font_size\":24,\"font_weight\":\"regular|bold\",\"font_weight_value\":700,\"font_family_hint\":\"sans\","
        "\"color\":\"#RRGGBB\",\"recognized_text\":\"...\",\"confidence\":0.0}],"
        "\"text_lines\":[{\"line_id\":\"...\",\"slot_id\":\"...\",\"order\":1,\"text\":\"...\","
        "\"bbox\":{\"x\":0,\"y\":0,\"w\":0,\"h\":0},\"align\":\"left|center|right\","
        "\"font_size\":24,\"font_weight\":\"regular|bold\",\"font_weight_value\":700,\"font_family_hint\":\"sans\",\"color\":\"#RRGGBB\",\"confidence\":0.0}],"
        "\"text_spans\":[{\"span_id\":\"...\",\"line_id\":\"...\",\"slot_id\":\"...\",\"order\":1,\"text\":\"...\","
        "\"char_start\":0,\"char_end\":2,"
        "\"bbox\":{\"x\":0,\"y\":0,\"w\":0,\"h\":0},\"font_size\":24,\"font_weight\":\"regular|bold\",\"font_weight_value\":700,\"font_family_hint\":\"sans\","
        "\"color\":\"#RRGGBB\",\"confidence\":0.0}]}\n"
        "Whitespace fidelity is mandatory: preserve spaces and punctuation exactly in span text.\n"
        "Do not trim span text.\n"
        "For each line, concatenating span.text by order must equal line.text exactly (character-for-character).\n"
        "char_start/char_end are 0-based, end-exclusive indexes into line.text.\n"
        "Coordinates must be normalized 0..1 relative to the full slide.\n\n"
        f"Schema slots:\n{json.dumps(slots, ensure_ascii=False)}"
    )


def extract_responses_output_text(resp_json: Dict[str, Any]) -> str:
    if isinstance(resp_json.get("output_text"), str) and resp_json.get("output_text"):
        return str(resp_json["output_text"])
    chunks: List[str] = []
    for item in resp_json.get("output", []) or []:
        if not isinstance(item, dict):
            continue
        for content in item.get("content", []) or []:
            if not isinstance(content, dict):
                continue
            ctype = str(content.get("type") or "")
            if ctype in {"output_text", "text"} and content.get("text"):
                chunks.append(str(content["text"]))
    if chunks:
        return "\n".join(chunks)
    raise ValueError(f"cannot extract output text from responses payload keys={list(resp_json.keys())}")


def infer_layout_with_vlm(
    image_path: Path,
    slide: Dict[str, Any],
    env: Dict[str, str],
    log: RunLogger | None = None,
    *,
    provider: str = "qwen",
    model_override: str = "",
    endpoint_override: str = "",
    api_key_override: str = "",
) -> Dict[str, Any]:
    prompt = build_vlm_prompt(slide)
    image_b64 = base64.b64encode(image_path.read_bytes()).decode("ascii")
    chosen = str(provider or "qwen").strip().lower()
    if chosen not in {"qwen", "gpt54"}:
        chosen = "qwen"
    if chosen == "gpt54":
        api_key = (
            api_key_override.strip()
            or env.get("AZURE_VLM_GPT54_API_KEY", "").strip()
            or env.get("AZURE_IMAGE_OPENAI_API_KEY", "").strip()
        )
        endpoint = (
            endpoint_override.strip()
            or env.get("AZURE_VLM_GPT54_ENDPOINT", "").strip()
            or env.get("AZURE_VLM_RESPONSES_ENDPOINT", "").strip()
        )
        model = model_override.strip() or env.get("AZURE_VLM_GPT54_MODEL", "").strip() or "gpt-5.4"
        if not endpoint or not api_key:
            if log:
                reason = "missing_AZURE_VLM_GPT54_ENDPOINT" if not endpoint else "missing_AZURE_VLM_GPT54_API_KEY"
                log.log("vlm_skipped", provider=chosen, reason=reason)
            return default_layout(slide)
        payload = {
            "model": model,
            "input": [
                {
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": prompt},
                        {"type": "input_image", "image_url": f"data:image/png;base64,{image_b64}"},
                    ],
                }
            ],
            "temperature": 0,
        }
        try:
            started = time.monotonic()
            if log:
                log.log(
                    "vlm_request",
                    provider=chosen,
                    endpoint=endpoint,
                    model=model,
                    image_path=str(image_path),
                    image_bytes=image_path.stat().st_size,
                    slot_count=len(slide.get("text_slots", [])),
                )
            with httpx.Client(timeout=120, follow_redirects=True) as client:
                resp = client.post(endpoint, headers={"api-key": api_key, "Content-Type": "application/json"}, json=payload)
                if log:
                    log.log(
                        "vlm_response",
                        provider=chosen,
                        endpoint=endpoint,
                        status_code=resp.status_code,
                        elapsed_ms=int((time.monotonic() - started) * 1000),
                        content_type=resp.headers.get("content-type", ""),
                        response_preview=(resp.text or "")[:1200],
                    )
                if resp.status_code >= 400:
                    return {**default_layout(slide), "vlm_error": f"http_{resp.status_code}: {resp.text[:500]}"}
                content = extract_responses_output_text(resp.json())
            parsed = parse_json_object(str(content))
            layout = normalize_layout(parsed, slide)
            layout["text_validation"] = build_text_validation(layout, slide)
            if log:
                log.log(
                    "vlm_layout_parsed",
                    provider=chosen,
                    slot_count=len(layout.get("text_slots", [])),
                    line_count=len(layout.get("text_lines", [])),
                    span_count=len(layout.get("text_spans", [])),
                    source=layout.get("source"),
                )
                log.log("vlm_text_validation", provider=chosen, items=layout["text_validation"])
            return layout
        except Exception as exc:
            if log:
                log.log("vlm_exception", provider=chosen, exception_type=type(exc).__name__, error=str(exc))
            return {**default_layout(slide), "vlm_error": str(exc)}

    # qwen/dashscope-compatible path
    api_key = api_key_override.strip() or (env.get("DASHSCOPE_API_KEY") or env.get("QWEN_API_KEY") or "").strip()
    base_url = (endpoint_override.strip() or env.get("VISION_BASE_URL", "https://dashscope.aliyuncs.com/compatible-mode/v1")).rstrip("/")
    model = model_override.strip() or env.get("VISION_MODEL", "qwen3-vl-plus").strip() or "qwen3-vl-plus"
    if not api_key:
        if log:
            log.log("vlm_skipped", provider=chosen, reason="missing_DASHSCOPE_API_KEY_or_QWEN_API_KEY")
        return default_layout(slide)
    payload = {
        "model": model,
        "messages": [
            {
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_b64}"}},
                    {"type": "text", "text": prompt},
                ],
            }
        ],
        "temperature": 0,
    }
    try:
        started = time.monotonic()
        if log:
            log.log(
                "vlm_request",
                provider=chosen,
                endpoint=f"{base_url}/chat/completions",
                model=model,
                image_path=str(image_path),
                image_bytes=image_path.stat().st_size,
                slot_count=len(slide.get("text_slots", [])),
            )
        with httpx.Client(timeout=120, follow_redirects=True) as client:
            resp = client.post(
                f"{base_url}/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json=payload,
            )
            if log:
                log.log(
                    "vlm_response",
                    provider=chosen,
                    endpoint=f"{base_url}/chat/completions",
                    status_code=resp.status_code,
                    elapsed_ms=int((time.monotonic() - started) * 1000),
                    content_type=resp.headers.get("content-type", ""),
                    response_preview=(resp.text or "")[:1200],
                )
            if resp.status_code >= 400:
                return {**default_layout(slide), "vlm_error": f"http_{resp.status_code}: {resp.text[:500]}"}
            content = resp.json()["choices"][0]["message"]["content"]
        parsed = parse_json_object(str(content))
        layout = normalize_layout(parsed, slide)
        layout["text_validation"] = build_text_validation(layout, slide)
        if log:
            log.log(
                "vlm_layout_parsed",
                provider=chosen,
                slot_count=len(layout.get("text_slots", [])),
                line_count=len(layout.get("text_lines", [])),
                span_count=len(layout.get("text_spans", [])),
                source=layout.get("source"),
            )
            log.log("vlm_text_validation", provider=chosen, items=layout["text_validation"])
        return layout
    except Exception as exc:
        if log:
            log.log("vlm_exception", provider=chosen, exception_type=type(exc).__name__, error=str(exc))
        return {**default_layout(slide), "vlm_error": str(exc)}


def parse_json_object(text: str) -> Dict[str, Any]:
    match = re.search(r"\{[\s\S]*\}", text)
    if not match:
        raise ValueError("no JSON object found")
    return json.loads(match.group(0))


def normalize_layout(layout: Dict[str, Any], slide: Dict[str, Any]) -> Dict[str, Any]:
    fallback_by_id = {slot["id"]: slot for slot in default_layout(slide)["text_slots"]}
    schema_text_by_id = {str(slot.get("id")): str(slot.get("text") or "") for slot in slide.get("text_slots", [])}
    normalized = {"source": "vlm", "meta": {"version": "2.0"}, "text_slots": [], "text_lines": [], "text_spans": []}
    for item in layout.get("text_slots", []):
        slot_id = str(item.get("id") or "").strip()
        if not slot_id or slot_id not in fallback_by_id:
            continue
        fallback = fallback_by_id[slot_id]
        bbox = item.get("bbox") if isinstance(item.get("bbox"), dict) else {}
        merged = dict(fallback)
        merged["bbox"] = {
            "x": clamp_float(bbox.get("x"), fallback["bbox"]["x"]),
            "y": clamp_float(bbox.get("y"), fallback["bbox"]["y"]),
            "w": clamp_float(bbox.get("w"), fallback["bbox"]["w"]),
            "h": clamp_float(bbox.get("h"), fallback["bbox"]["h"]),
        }
        align = str(item.get("align") or fallback.get("align") or "left").strip().lower()
        merged["align"] = align if align in {"left", "center", "right"} else str(fallback.get("align") or "left")
        fallback_weight_token, fallback_weight_value = parse_font_weight(fallback.get("font_weight"), "regular")
        parsed_weight_token, parsed_weight_value = parse_font_weight(item.get("font_weight"), fallback_weight_token)
        if item.get("font_weight_value") is not None:
            parsed_weight_token, parsed_weight_value = parse_font_weight(item.get("font_weight_value"), parsed_weight_token)
        merged["font_weight"] = "bold" if parsed_weight_token == "bold" or parsed_weight_value >= 600 else "regular"
        merged["font_weight_value"] = parsed_weight_value
        merged["font_family_hint"] = normalize_font_family_hint(item.get("font_family_hint"), normalize_font_family_hint(fallback.get("font_family_hint"), "sans"))
        merged["font_size"] = normalize_font_size(
            item.get("font_size"),
            fallback.get("font_size", 20),
            min_scale=0.75,
            max_scale=1.35,
        )
        inferred_color = normalize_hex_color(item.get("color"), str(fallback.get("color") or "#ffffff"))
        # Keep inferred color by default; only fallback when confidence is very low.
        if clamp_confidence(item.get("confidence"), 0.0) < 0.25:
            merged["color"] = normalize_hex_color(fallback.get("color"), inferred_color)
        else:
            merged["color"] = inferred_color
        merged["recognized_text"] = str(item.get("recognized_text") or item.get("text") or "").strip()
        merged["confidence"] = clamp_confidence(item.get("confidence"), 0.0)
        normalized["text_slots"].append(merged)
    if len(normalized["text_slots"]) != len(fallback_by_id):
        seen = {s["id"] for s in normalized["text_slots"]}
        for slot_id, fallback in fallback_by_id.items():
            if slot_id not in seen:
                normalized["text_slots"].append(fallback)
    line_index_by_slot: Dict[str, int] = {}
    raw_lines = layout.get("text_lines") or layout.get("lines") or []
    for line in raw_lines:
        slot_id = str(line.get("slot_id") or line.get("id") or "").strip()
        if not slot_id or slot_id not in fallback_by_id:
            continue
        fallback = fallback_by_id[slot_id]
        bbox = line.get("bbox") if isinstance(line.get("bbox"), dict) else {}
        entry = {
            "line_id": str(line.get("line_id") or "").strip(),
            "order": int(line.get("order") or 0) if str(line.get("order") or "").strip().isdigit() else 0,
            "slot_id": slot_id,
            "text": str(line.get("text") or line.get("recognized_text") or "").strip(),
            "bbox": {
                "x": clamp_float(bbox.get("x"), fallback["bbox"]["x"]),
                "y": clamp_float(bbox.get("y"), fallback["bbox"]["y"]),
                "w": clamp_float(bbox.get("w"), fallback["bbox"]["w"]),
                "h": clamp_float(bbox.get("h"), fallback["bbox"]["h"]),
            },
            "align": str(line.get("align") or fallback.get("align") or "left").strip().lower(),
            "font_size": normalize_font_size(line.get("font_size"), fallback.get("font_size", 20), min_scale=0.7, max_scale=1.5),
            "font_weight": "regular",
            "font_weight_value": 400,
            "font_family_hint": normalize_font_family_hint(line.get("font_family_hint"), normalize_font_family_hint(fallback.get("font_family_hint"), "sans")),
            "color": normalize_hex_color(line.get("color"), str(fallback.get("color") or "#ffffff")),
            "confidence": clamp_confidence(line.get("confidence"), 0.0),
        }
        line_weight_token, line_weight_value = parse_font_weight(line.get("font_weight"), str(fallback.get("font_weight") or "regular"))
        if line.get("font_weight_value") is not None:
            line_weight_token, line_weight_value = parse_font_weight(line.get("font_weight_value"), line_weight_token)
        entry["font_weight"] = "bold" if line_weight_token == "bold" or line_weight_value >= 600 else "regular"
        entry["font_weight_value"] = line_weight_value
        if not entry["text"]:
            continue
        if entry["align"] not in {"left", "center", "right"}:
            entry["align"] = str(fallback.get("align") or "left")
        if not entry["line_id"]:
            line_index_by_slot[slot_id] = int(line_index_by_slot.get(slot_id, 0)) + 1
            entry["line_id"] = f"{slot_id}_l{line_index_by_slot[slot_id]}"
        if entry["order"] <= 0:
            order_match = re.search(r"_l(\d+)$", entry["line_id"])
            if order_match:
                entry["order"] = int(order_match.group(1))
            else:
                entry["order"] = int(line_index_by_slot.get(slot_id, 1))
        normalized["text_lines"].append(entry)
    if not normalized["text_lines"]:
        for slot in normalized["text_slots"]:
            line_id = f"{slot['id']}_l1"
            normalized["text_lines"].append(
                {
                    "line_id": line_id,
                    "order": 1,
                    "slot_id": slot["id"],
                    "text": str(slot.get("recognized_text") or "").strip() or schema_text_by_id.get(str(slot["id"]), ""),
                    "bbox": dict(slot.get("bbox") or fallback_by_id[slot["id"]]["bbox"]),
                    "align": slot.get("align", "left"),
                    "font_size": slot.get("font_size", 20),
                    "font_weight": slot.get("font_weight", "regular"),
                    "font_weight_value": slot.get("font_weight_value", 400),
                    "font_family_hint": slot.get("font_family_hint", "sans"),
                    "color": slot.get("color", "#ffffff"),
                    "confidence": slot.get("confidence", 0.0),
                }
            )
    normalized["text_lines"].sort(
        key=lambda x: (
            x.get("slot_id", ""),
            int(x.get("order") or 0),
            float((x.get("bbox") or {}).get("y", 0)),
        )
    )
    lines_by_id: Dict[str, Dict[str, Any]] = {str(line.get("line_id") or ""): line for line in normalized["text_lines"]}
    raw_spans = layout.get("text_spans") or layout.get("spans") or []
    for span in raw_spans:
        slot_id = str(span.get("slot_id") or "").strip()
        line_id = str(span.get("line_id") or "").strip()
        if line_id and line_id in lines_by_id:
            slot_id = slot_id or str(lines_by_id[line_id].get("slot_id") or "")
        if not slot_id or slot_id not in fallback_by_id:
            continue
        fallback = fallback_by_id[slot_id]
        if not line_id:
            line_id = f"{slot_id}_l1"
        line_ref = lines_by_id.get(line_id) or {}
        bbox = span.get("bbox") if isinstance(span.get("bbox"), dict) else {}
        text = str(span.get("text") or "")
        start_idx = parse_optional_nonneg_int(span.get("char_start"))
        end_idx = parse_optional_nonneg_int(span.get("char_end"))
        if text == "" and not (start_idx is not None and end_idx is not None and end_idx > start_idx):
            continue
        entry = {
            "span_id": str(span.get("span_id") or "").strip() or f"{line_id}_s1",
            "line_id": line_id,
            "slot_id": slot_id,
            "order": int(span.get("order") or 0) if str(span.get("order") or "").strip().isdigit() else 0,
            "text": text,
            "char_start": start_idx,
            "char_end": end_idx,
            "bbox": {
                "x": clamp_float(bbox.get("x"), (line_ref.get("bbox") or fallback["bbox"]).get("x", fallback["bbox"]["x"])),
                "y": clamp_float(bbox.get("y"), (line_ref.get("bbox") or fallback["bbox"]).get("y", fallback["bbox"]["y"])),
                "w": clamp_float(bbox.get("w"), (line_ref.get("bbox") or fallback["bbox"]).get("w", fallback["bbox"]["w"])),
                "h": clamp_float(bbox.get("h"), (line_ref.get("bbox") or fallback["bbox"]).get("h", fallback["bbox"]["h"])),
            },
            "font_size": normalize_font_size(span.get("font_size"), line_ref.get("font_size") or fallback.get("font_size", 20), min_scale=0.7, max_scale=1.5),
            "font_weight": "regular",
            "font_weight_value": 400,
            "font_family_hint": normalize_font_family_hint(
                span.get("font_family_hint"),
                normalize_font_family_hint(line_ref.get("font_family_hint"), normalize_font_family_hint(fallback.get("font_family_hint"), "sans")),
            ),
            "color": normalize_hex_color(span.get("color"), str(line_ref.get("color") or fallback.get("color") or "#ffffff")),
            "confidence": clamp_confidence(span.get("confidence"), line_ref.get("confidence") or 0.0),
        }
        span_weight_token, span_weight_value = parse_font_weight(span.get("font_weight"), str(line_ref.get("font_weight") or fallback.get("font_weight") or "regular"))
        if span.get("font_weight_value") is not None:
            span_weight_token, span_weight_value = parse_font_weight(span.get("font_weight_value"), span_weight_token)
        entry["font_weight"] = "bold" if span_weight_token == "bold" or span_weight_value >= 600 else "regular"
        entry["font_weight_value"] = span_weight_value
        if entry["order"] <= 0:
            entry["order"] = int(len([s for s in normalized["text_spans"] if str(s.get("line_id")) == line_id]) + 1)
        normalized["text_spans"].append(entry)
    if not normalized["text_spans"]:
        for line in normalized["text_lines"]:
            line_id = str(line.get("line_id") or "")
            line_text = str(line.get("text") or "")
            normalized["text_spans"].append(
                {
                    "span_id": f"{line_id}_s1",
                    "line_id": line_id,
                    "slot_id": str(line.get("slot_id") or ""),
                    "order": 1,
                    "text": line_text,
                    "char_start": 0,
                    "char_end": len(line_text),
                    "bbox": dict(line.get("bbox") or {}),
                    "font_size": line.get("font_size", 20),
                    "font_weight": line.get("font_weight", "regular"),
                    "font_weight_value": line.get("font_weight_value", 400),
                    "font_family_hint": line.get("font_family_hint", "sans"),
                    "color": line.get("color", "#ffffff"),
                    "confidence": line.get("confidence", 0.0),
                }
            )
    normalized["text_spans"].sort(
        key=lambda x: (
            x.get("slot_id", ""),
            str(x.get("line_id") or ""),
            int(x.get("order") or 0),
            float((x.get("bbox") or {}).get("x", 0)),
        )
    )
    normalized["text_spans"] = reconcile_spans_with_lines(normalized["text_lines"], normalized["text_spans"])
    return normalized


def reconcile_spans_with_lines(lines: List[Dict[str, Any]], spans: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    lines_by_id: Dict[str, Dict[str, Any]] = {str(line.get("line_id") or ""): line for line in lines}
    spans_by_line: Dict[str, List[Dict[str, Any]]] = {}
    for span in spans:
        lid = str(span.get("line_id") or "")
        if lid:
            spans_by_line.setdefault(lid, []).append(span)
    rebuilt: List[Dict[str, Any]] = []
    for line in lines:
        line_id = str(line.get("line_id") or "")
        line_text = str(line.get("text") or "")
        line_spans = sorted(
            spans_by_line.get(line_id, []),
            key=lambda s: (
                int(s.get("order") or 0),
                parse_optional_nonneg_int(s.get("char_start")) if parse_optional_nonneg_int(s.get("char_start")) is not None else 10**9,
                float((s.get("bbox") or {}).get("x", 0)),
            ),
        )
        if not line_spans:
            rebuilt.append(
                {
                    "span_id": f"{line_id}_s1",
                    "line_id": line_id,
                    "slot_id": str(line.get("slot_id") or ""),
                    "order": 1,
                    "text": line_text,
                    "char_start": 0,
                    "char_end": len(line_text),
                    "bbox": dict(line.get("bbox") or {}),
                    "font_size": line.get("font_size", 20),
                    "font_weight": line.get("font_weight", "regular"),
                    "font_weight_value": line.get("font_weight_value", 400),
                    "font_family_hint": line.get("font_family_hint", "sans"),
                    "color": line.get("color", "#ffffff"),
                    "confidence": line.get("confidence", 0.0),
                }
            )
            continue

        indexed_ok = all(parse_optional_nonneg_int(s.get("char_start")) is not None and parse_optional_nonneg_int(s.get("char_end")) is not None for s in line_spans)
        if indexed_ok and line_text != "":
            cursor = 0
            out_line: List[Dict[str, Any]] = []
            out_order = 1
            for src in line_spans:
                s = max(0, min(len(line_text), int(parse_optional_nonneg_int(src.get("char_start")) or 0)))
                e = max(s, min(len(line_text), int(parse_optional_nonneg_int(src.get("char_end")) or s)))
                if s > cursor:
                    gap_txt = line_text[cursor:s]
                    if gap_txt:
                        out_line.append(
                            {
                                **src,
                                "span_id": f"{line_id}_auto_gap_{out_order}",
                                "order": out_order,
                                "text": gap_txt,
                                "char_start": cursor,
                                "char_end": s,
                            }
                        )
                        out_order += 1
                txt = line_text[s:e]
                if txt:
                    out_line.append(
                        {
                            **src,
                            "order": out_order,
                            "text": txt,
                            "char_start": s,
                            "char_end": e,
                        }
                    )
                    out_order += 1
                cursor = max(cursor, e)
            if cursor < len(line_text):
                tail_txt = line_text[cursor:]
                if tail_txt:
                    ref = out_line[-1] if out_line else line_spans[-1]
                    out_line.append(
                        {
                            **ref,
                            "span_id": f"{line_id}_auto_tail_{out_order}",
                            "order": out_order,
                            "text": tail_txt,
                            "char_start": cursor,
                            "char_end": len(line_text),
                        }
                    )
            if "".join(str(s.get("text") or "") for s in out_line) == line_text:
                rebuilt.extend(out_line)
                continue

        # text-only fallback check
        joined = "".join(str(s.get("text") or "") for s in line_spans)
        if joined != line_text:
            rebuilt.append(
                {
                    "span_id": f"{line_id}_s1",
                    "line_id": line_id,
                    "slot_id": str(line.get("slot_id") or ""),
                    "order": 1,
                    "text": line_text,
                    "char_start": 0,
                    "char_end": len(line_text),
                    "bbox": dict(line.get("bbox") or {}),
                    "font_size": line.get("font_size", 20),
                    "font_weight": line.get("font_weight", "regular"),
                    "font_weight_value": line.get("font_weight_value", 400),
                    "font_family_hint": line.get("font_family_hint", "sans"),
                    "color": line.get("color", "#ffffff"),
                    "confidence": line.get("confidence", 0.0),
                }
            )
            continue
        rebuilt.extend(line_spans)

    rebuilt.sort(
        key=lambda x: (
            str(x.get("slot_id") or ""),
            str(x.get("line_id") or ""),
            int(x.get("order") or 0),
            float((x.get("bbox") or {}).get("x", 0)),
        )
    )
    return rebuilt


def normalize_text_for_compare(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "").strip()).lower()


def build_text_validation(layout: Dict[str, Any], slide: Dict[str, Any]) -> List[Dict[str, Any]]:
    schema_by_id = {str(s.get("id")): str(s.get("text") or "") for s in slide.get("text_slots", [])}
    lines_by_slot: Dict[str, List[Dict[str, Any]]] = {}
    for line in layout.get("text_lines", []):
        slot_id = str(line.get("slot_id") or "")
        if not slot_id:
            continue
        lines_by_slot.setdefault(slot_id, []).append(line)
    out: List[Dict[str, Any]] = []
    for slot_id, expected in schema_by_id.items():
        ordered_lines = sorted(
            lines_by_slot.get(slot_id, []),
            key=lambda x: (
                int(x.get("order") or 0),
                float((x.get("bbox") or {}).get("y", 0)),
            ),
        )
        recognized = "\n".join([str(x.get("text") or "").strip() for x in ordered_lines if str(x.get("text") or "").strip()]).strip()
        if not recognized:
            for slot in layout.get("text_slots", []):
                if str(slot.get("id")) == slot_id:
                    recognized = str(slot.get("recognized_text") or "").strip()
                    break
        out.append(
            {
                "slot_id": slot_id,
                "expected_text": expected,
                "recognized_text": recognized,
                "normalized_match": normalize_text_for_compare(expected) == normalize_text_for_compare(recognized),
            }
        )
    return out


def clamp_float(value: Any, default: float) -> float:
    try:
        parsed = float(value)
    except Exception:
        parsed = float(default)
    return max(0.0, min(1.0, parsed))


def clamp_confidence(value: Any, default: float) -> float:
    try:
        parsed = float(value)
    except Exception:
        parsed = float(default)
    return max(0.0, min(1.0, parsed))


def parse_optional_nonneg_int(value: Any) -> int | None:
    try:
        out = int(value)
    except Exception:
        return None
    if out < 0:
        return None
    return out


def normalize_font_size(value: Any, default: Any, *, min_scale: float = 0.5, max_scale: float = 2.0) -> int:
    try:
        parsed = int(round(float(value)))
    except Exception:
        try:
            parsed = int(round(float(default)))
        except Exception:
            parsed = 20
    try:
        base = float(default)
    except Exception:
        base = 20.0
    lower = max(8, int(round(base * min_scale)))
    upper = min(200, int(round(base * max_scale)))
    if lower > upper:
        lower, upper = upper, lower
    return max(lower, min(upper, parsed))


def normalize_hex_color(value: Any, default: str) -> str:
    raw = str(value or "").strip()
    if not raw:
        raw = str(default or "").strip()
    if not raw:
        return "#ffffff"
    rgb_match = re.match(r"rgba?\((\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})", raw, flags=re.IGNORECASE)
    if rgb_match:
        r = max(0, min(255, int(rgb_match.group(1))))
        g = max(0, min(255, int(rgb_match.group(2))))
        b = max(0, min(255, int(rgb_match.group(3))))
        return f"#{r:02x}{g:02x}{b:02x}"
    if raw.startswith("#"):
        hex_part = raw[1:]
        if re.fullmatch(r"[0-9a-fA-F]{3}", hex_part):
            return "#" + "".join(ch * 2 for ch in hex_part).lower()
        if re.fullmatch(r"[0-9a-fA-F]{6}", hex_part):
            return "#" + hex_part.lower()
    if re.fullmatch(r"[0-9a-fA-F]{6}", raw):
        return "#" + raw.lower()
    return normalize_hex_color(default, "#ffffff") if raw != str(default or "").strip() else "#ffffff"


def build_mask(layout: Dict[str, Any], out_path: Path) -> None:
    mask = Image.new("RGBA", (SLIDE_W, SLIDE_H), (255, 255, 255, 0))
    draw = ImageDraw.Draw(mask)
    regions = layout.get("text_spans") or layout.get("text_lines") or layout.get("text_slots") or []
    for item in regions:
        bbox = item.get("bbox") or {}
        x = int(float(bbox.get("x", 0)) * SLIDE_W)
        y = int(float(bbox.get("y", 0)) * SLIDE_H)
        w = int(float(bbox.get("w", 0)) * SLIDE_W)
        h = int(float(bbox.get("h", 0)) * SLIDE_H)
        if w <= 1 or h <= 1:
            continue

        # Stable rectangular mask baseline.
        pad_x = max(2, min(8, int(w * 0.03)))
        pad_y = max(2, min(6, int(h * 0.08)))
        x0 = max(0, x - pad_x)
        y0 = max(0, y - pad_y)
        x1 = min(SLIDE_W, x + w + pad_x)
        y1 = min(SLIDE_H, y + h + pad_y)
        if x1 - x0 <= 1 or y1 - y0 <= 1:
            continue
        radius = max(2, min(10, int(min(x1 - x0, y1 - y0) * 0.18)))
        draw.rounded_rectangle((x0, y0, x1, y1), radius=radius, fill=(255, 255, 255, 255))
    mask.save(out_path)


def image_edit_prompt(strategy: str = "masked") -> str:
    if str(strategy).strip().lower() == "prompt-only":
        return (
            "Remove all visible readable text from the entire image while preserving layout and visual design. "
            "Keep composition, objects, icons, borders, lighting, and background style consistent. "
            "Do not add any new text, letters, numbers, labels, watermarks, logo text, or pseudo text."
        )
    return (
        "Remove only text inside masked regions and fill those areas as natural background. "
        "Do not add any new text, letters, numbers, labels, watermarks, logo text, or pseudo text."
    )


def dry_remove_text(image_path: Path, mask_path: Path, out_path: Path) -> None:
    img = Image.open(image_path).convert("RGB")
    mask = Image.open(mask_path).convert("L")
    blurred = img.filter(ImageFilter.GaussianBlur(radius=22))
    cleaned = Image.composite(blurred, img, mask)
    cleaned.save(out_path)


def local_inpaint_text(image_path: Path, mask_path: Path, out_path: Path, log: RunLogger | None = None) -> None:
    """
    Local text removal only: no model edit call.
    Use component-wise inpainting to reduce collateral damage on nearby UI structures.
    """
    if cv2 is None:
        if log:
            log.log("image_edit_local_fallback", reason="cv2_unavailable", method="gaussian_blur_composite")
        dry_remove_text(image_path, mask_path, out_path)
        return

    img_bytes = np.frombuffer(image_path.read_bytes(), dtype=np.uint8)
    mask_bytes = np.frombuffer(mask_path.read_bytes(), dtype=np.uint8)
    bgr = cv2.imdecode(img_bytes, cv2.IMREAD_COLOR)
    raw_mask = cv2.imdecode(mask_bytes, cv2.IMREAD_GRAYSCALE)
    if bgr is None or raw_mask is None:
        if log:
            log.log("image_edit_local_fallback", reason="decode_failed", method="gaussian_blur_composite")
        dry_remove_text(image_path, mask_path, out_path)
        return
    if raw_mask.shape[:2] != bgr.shape[:2]:
        raw_mask = cv2.resize(raw_mask, (bgr.shape[1], bgr.shape[0]), interpolation=cv2.INTER_LINEAR)

    _, bin_mask = cv2.threshold(raw_mask, 10, 255, cv2.THRESH_BINARY)
    kernel = cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (3, 3))
    bin_mask = cv2.morphologyEx(bin_mask, cv2.MORPH_CLOSE, kernel, iterations=1)
    bin_mask = cv2.morphologyEx(bin_mask, cv2.MORPH_OPEN, kernel, iterations=1)
    # Expand mask slightly to cover antialias/glow at glyph edges.
    inpaint_mask = cv2.dilate(bin_mask, cv2.getStructuringElement(cv2.MORPH_ELLIPSE, (5, 5)), iterations=1)

    out = bgr.copy()
    comp_count, labels, stats, _ = cv2.connectedComponentsWithStats(inpaint_mask, connectivity=8)
    edited_components = 0
    for idx in range(1, comp_count):
        x, y, w, h, area = stats[idx]
        if area <= 4 or w <= 1 or h <= 1:
            continue
        pad = max(8, min(56, int(max(w, h) * 0.35)))
        x0 = max(0, x - pad)
        y0 = max(0, y - pad)
        x1 = min(out.shape[1], x + w + pad)
        y1 = min(out.shape[0], y + h + pad)
        roi_img = out[y0:y1, x0:x1]
        roi_mask = inpaint_mask[y0:y1, x0:x1]
        radius = max(2, min(10, int(max(w, h) * 0.08)))
        pass1 = cv2.inpaint(roi_img, roi_mask, radius, cv2.INPAINT_TELEA)
        pass2 = cv2.inpaint(pass1, roi_mask, max(1, radius - 1), cv2.INPAINT_NS)
        keep = roi_mask > 0
        roi_out = roi_img.copy()
        roi_out[keep] = pass2[keep]
        out[y0:y1, x0:x1] = roi_out
        edited_components += 1

    if edited_components == 0:
        if log:
            log.log("image_edit_local_noop", reason="no_components")
        out = bgr
    ok, encoded = cv2.imencode(".png", out)
    if not ok:
        if log:
            log.log("image_edit_local_fallback", reason="encode_failed", method="gaussian_blur_composite")
        dry_remove_text(image_path, mask_path, out_path)
        return
    out_path.write_bytes(encoded.tobytes())
    if log:
        log.log(
            "image_edit_local_done",
            method="opencv_inpaint_componentwise",
            components=edited_components,
            path=str(out_path),
            bytes=out_path.stat().st_size,
        )


def make_pptx(schema: Dict[str, Any], layouts: Dict[str, Dict[str, Any]], background_paths: Dict[str, Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_schema in schema.get("slides", []):
        sid = str(slide_schema.get("id") or "")
        if sid not in background_paths or sid not in layouts:
            continue
        slide = prs.slides.add_slide(blank)
        bg_path = background_paths[sid]
        slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        layout = layouts[sid]
        schema_slots = {slot["id"]: slot for slot in slide_schema.get("text_slots", [])}
        lines_by_slot: Dict[str, List[Dict[str, Any]]] = {}
        for line in layout.get("text_lines", []):
            slot_id = str(line.get("slot_id") or "")
            if slot_id:
                lines_by_slot.setdefault(slot_id, []).append(line)
        spans_by_line: Dict[str, List[Dict[str, Any]]] = {}
        for span in layout.get("text_spans", []):
            line_id = str(span.get("line_id") or "")
            if line_id:
                spans_by_line.setdefault(line_id, []).append(span)
        for inferred in layout.get("text_slots", []):
            slot = schema_slots.get(inferred["id"])
            if not slot:
                continue
            bbox = inferred.get("bbox") or {}
            left = Inches(float(bbox.get("x", slot["x"])) * 13.333333)
            top = Inches(float(bbox.get("y", slot["y"])) * 7.5)
            width = Inches(float(bbox.get("w", slot["w"])) * 13.333333)
            height = Inches(float(bbox.get("h", slot["h"])) * 7.5)
            shape = slide.shapes.add_textbox(left, top, width, height)
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = False
            slot_lines = sorted(
                lines_by_slot.get(str(inferred["id"]), []),
                key=lambda x: (int(x.get("order") or 0), float((x.get("bbox") or {}).get("y", 0))),
            )
            if not slot_lines:
                slot_lines = [
                    {
                        "line_id": f"{inferred['id']}_l1",
                        "order": 1,
                        "slot_id": inferred["id"],
                        "text": str(slot.get("text", "")),
                        "align": inferred.get("align", slot.get("align", "left")),
                        "font_size": inferred.get("font_size", slot.get("font_size", 18)),
                        "font_weight": inferred.get("font_weight", slot.get("font_weight", "regular")),
                        "font_weight_value": inferred.get("font_weight_value", slot.get("font_weight_value", 400)),
                        "font_family_hint": inferred.get("font_family_hint", slot.get("font_family_hint", "sans")),
                        "color": inferred.get("color", slot.get("color", "#ffffff")),
                    }
                ]
            for idx, line in enumerate(slot_lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.alignment = ppt_align(line.get("align", inferred.get("align", slot.get("align", "left"))))
                line_id = str(line.get("line_id") or "")
                line_spans = sorted(
                    spans_by_line.get(line_id, []),
                    key=lambda x: (int(x.get("order") or 0), float((x.get("bbox") or {}).get("x", 0))),
                )
                if not line_spans:
                    line_spans = [
                        {
                            "text": str(line.get("text", "")),
                            "font_size": line.get("font_size", inferred.get("font_size", slot.get("font_size", 18))),
                            "font_weight": line.get("font_weight", inferred.get("font_weight", slot.get("font_weight", "regular"))),
                            "font_weight_value": line.get("font_weight_value", inferred.get("font_weight_value", slot.get("font_weight_value", 400))),
                            "font_family_hint": line.get("font_family_hint", inferred.get("font_family_hint", slot.get("font_family_hint", "sans"))),
                            "color": line.get("color", inferred.get("color", slot.get("color", "#ffffff"))),
                        }
                    ]
                for span in line_spans:
                    run = p.add_run()
                    run_text = str(span.get("text", ""))
                    run.text = run_text
                    run.font.size = Pt(float(span.get("font_size") or line.get("font_size") or inferred.get("font_size") or slot.get("font_size") or 18))
                    wt, wv = parse_font_weight(
                        span.get("font_weight") or line.get("font_weight") or inferred.get("font_weight") or slot.get("font_weight"),
                        "regular",
                    )
                    if span.get("font_weight_value") is not None:
                        wt, wv = parse_font_weight(span.get("font_weight_value"), wt)
                    elif line.get("font_weight_value") is not None:
                        wt, wv = parse_font_weight(line.get("font_weight_value"), wt)
                    run.font.bold = wt == "bold" or wv >= 600
                    fam = normalize_font_family_hint(
                        span.get("font_family_hint") or line.get("font_family_hint") or inferred.get("font_family_hint") or slot.get("font_family_hint"),
                        "sans",
                    )
                    run.font.name = str(font_targets(run_text, fam, wv, run.font.bold).get("ppt_name") or "Arial")
                    run.font.color.rgb = rgb_color(str(span.get("color") or line.get("color") or inferred.get("color") or slot.get("color") or "#ffffff"))
        for chart in slide_schema.get("chart_slots", []):
            add_editable_bar_chart(slide, chart)
    prs.save(out_path)


def ppt_align(value: str):
    token = str(value or "").lower()
    if token == "center":
        return PP_ALIGN.CENTER
    if token == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def rgb_color(value: str) -> RGBColor:
    raw = str(value or "#ffffff").strip().lstrip("#")
    if len(raw) != 6:
        raw = "ffffff"
    return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))


def add_editable_bar_chart(slide, chart: Dict[str, Any]) -> None:
    x, y, w, h = chart["x"], chart["y"], chart["w"], chart["h"]
    left = Inches(x * 13.333333)
    top = Inches(y * 7.5)
    width = Inches(w * 13.333333)
    height = Inches(h * 7.5)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.background()
    panel.line.color.rgb = RGBColor(103, 232, 249)
    labels = chart.get("labels", [])
    values = chart.get("values", [])
    if not values:
        return
    max_v = max(float(v) for v in values) or 1
    for idx, value in enumerate(values):
        bar_left = Inches((x + 0.05 + idx * (w / len(values))) * 13.333333)
        bar_h = (float(value) / max_v) * h * 0.62
        bar_top = Inches((y + h - 0.12 - bar_h) * 7.5)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            bar_left,
            bar_top,
            Inches((w / len(values)) * 13.333333 * 0.42),
            Inches(bar_h * 7.5),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(34, 211, 238)
        bar.line.fill.background()
        label = slide.shapes.add_textbox(bar_left, Inches((y + h - 0.08) * 7.5), Inches((w / len(values)) * 13.333333 * 0.7), Inches(0.25))
        label.text_frame.text = str(labels[idx]) if idx < len(labels) else ""


def render_overlay_preview(background_path: Path, slide: Dict[str, Any], layout: Dict[str, Any], out_path: Path) -> None:
    img = Image.open(background_path).convert("RGB").resize((SLIDE_W, SLIDE_H))
    draw = ImageDraw.Draw(img, "RGBA")
    slots = {slot["id"]: slot for slot in slide.get("text_slots", [])}
    lines_by_slot: Dict[str, List[Dict[str, Any]]] = {}
    for line in layout.get("text_lines", []):
        slot_id = str(line.get("slot_id") or "")
        if slot_id:
            lines_by_slot.setdefault(slot_id, []).append(line)
    spans_by_line: Dict[str, List[Dict[str, Any]]] = {}
    for span in layout.get("text_spans", []):
        line_id = str(span.get("line_id") or "")
        if line_id:
            spans_by_line.setdefault(line_id, []).append(span)
    for inferred in layout.get("text_slots", []):
        slot = slots.get(inferred["id"])
        if not slot:
            continue
        slot_lines = sorted(
            lines_by_slot.get(str(inferred["id"]), []),
            key=lambda x: (int(x.get("order") or 0), float((x.get("bbox") or {}).get("y", 0))),
        )
        if slot_lines:
            for line in slot_lines:
                line_bbox = line.get("bbox") or {}
                line_text = str(line.get("text") or "").strip()
                if not line_text:
                    continue
                line_slot = {
                    **slot,
                    "x": line_bbox.get("x", inferred.get("bbox", {}).get("x", slot["x"])),
                    "y": line_bbox.get("y", inferred.get("bbox", {}).get("y", slot["y"])),
                    "w": line_bbox.get("w", inferred.get("bbox", {}).get("w", slot["w"])),
                    "h": line_bbox.get("h", inferred.get("bbox", {}).get("h", slot["h"])),
                    "font_size": line.get("font_size") or inferred.get("font_size") or slot.get("font_size"),
                    "font_weight": line.get("font_weight") or inferred.get("font_weight") or slot.get("font_weight"),
                    "font_weight_value": line.get("font_weight_value") or inferred.get("font_weight_value") or slot.get("font_weight_value", 400),
                    "font_family_hint": line.get("font_family_hint") or inferred.get("font_family_hint") or slot.get("font_family_hint", "sans"),
                    "color": line.get("color") or inferred.get("color") or slot.get("color"),
                    "align": line.get("align") or inferred.get("align") or slot.get("align"),
                    "role": slot.get("role"),
                    "text": line_text,
                }
                line_id = str(line.get("line_id") or "")
                line_spans = spans_by_line.get(line_id, [])
                if line_spans:
                    draw_line_spans(draw, line_slot, line_spans, scale=1.0)
                else:
                    # Line-level render: disable wrapping by using single-line slot height/width.
                    draw_slot_text(draw, line_slot, line_text, scale=1.0, force_no_wrap=True)
            continue
        bbox = inferred.get("bbox") or {}
        draw_slot_text(
            draw,
            {
                **slot,
                "x": bbox.get("x", slot["x"]),
                "y": bbox.get("y", slot["y"]),
                "w": bbox.get("w", slot["w"]),
                "h": bbox.get("h", slot["h"]),
                "font_size": inferred.get("font_size") or slot.get("font_size"),
                "font_weight": inferred.get("font_weight") or slot.get("font_weight"),
                "font_weight_value": inferred.get("font_weight_value") or slot.get("font_weight_value", 400),
                "font_family_hint": inferred.get("font_family_hint") or slot.get("font_family_hint", "sans"),
                "color": inferred.get("color") or slot.get("color"),
                "align": inferred.get("align") or slot.get("align"),
            },
            str(slot.get("text", "")),
        )
    for chart in slide.get("chart_slots", []):
        draw_chart_preview(draw, chart)
    img.save(out_path)


def build_html(schema: Dict[str, Any], run_dir: Path, manifest: Dict[str, Any]) -> None:
    rows = []
    errors = list(manifest.get("errors") or [])
    stage = str(manifest.get("stage") or "full")
    error_html = ""
    if errors:
        error_items = "".join(
            "<li><strong>{slide}</strong> / {stage}: {error}</li>".format(
                slide=str(item.get("slide", "")),
                stage=str(item.get("stage", "")),
                error=str(item.get("error", "")),
            )
            for item in errors
        )
        error_html = f"""
        <section class="errors">
          <h2>API Errors</h2>
          <p>This run contains fallback images. Do not evaluate it as a real image model result.</p>
          <ul>{error_items}</ul>
        </section>
        """
    generated_slide_ids = set((manifest.get("slides") or {}).keys())
    for slide in schema.get("slides", []):
        sid = slide["id"]
        if sid not in generated_slide_ids:
            continue
        if stage == "generate-only":
            rows.append(
                f"""
                <section class="slide-block">
                  <h2>{sid}: {slide.get('type', '')}</h2>
                  <div class="grid single">
                    <figure><img src="{sid}_A.png"><figcaption>A visual draft with burned-in text</figcaption></figure>
                  </div>
                </section>
                """
            )
            continue
        rows.append(
            f"""
            <section class="slide-block">
              <h2>{sid}: {slide.get('type', '')}</h2>
              <div class="grid">
                <figure><img src="{sid}_A.png"><figcaption>A visual draft with burned-in text</figcaption></figure>
                <figure><img src="{sid}_mask.png"><figcaption>Generated text mask</figcaption></figure>
                <figure><img src="{sid}_B.png"><figcaption>B edited no-text background</figcaption></figure>
                <figure><img src="{sid}_C_preview.png"><figcaption>C preview: B + editable overlay positions</figcaption></figure>
              </div>
              <details><summary>Inferred layout JSON</summary><pre>{json.dumps(manifest.get('layouts', {}).get(sid, {}), ensure_ascii=False, indent=2)}</pre></details>
            </section>
            """
        )
    html = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <title>PPT Image-Native Visual PoC</title>
  <style>
    body {{ margin: 0; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; background: #0b1020; color: #e5e7eb; }}
    header {{ padding: 28px 36px; border-bottom: 1px solid rgba(255,255,255,.12); }}
    h1 {{ margin: 0 0 10px; font-size: 24px; }}
    a {{ color: #67e8f9; }}
    .meta {{ color: #94a3b8; font-size: 14px; }}
    .slide-block {{ padding: 28px 36px 38px; border-bottom: 1px solid rgba(255,255,255,.1); }}
    .errors {{ margin: 24px 36px 0; padding: 18px 20px; border: 1px solid rgba(248,113,113,.65); background: rgba(127,29,29,.35); border-radius: 8px; }}
    .errors h2 {{ margin: 0 0 8px; color: #fecaca; }}
    .errors p {{ margin: 0 0 10px; color: #fee2e2; }}
    .errors li {{ margin: 6px 0; color: #fecaca; }}
    h2 {{ margin: 0 0 18px; font-size: 18px; }}
    .grid {{ display: grid; grid-template-columns: repeat(2, minmax(320px, 1fr)); gap: 18px; }}
    .grid.single {{ grid-template-columns: minmax(320px, 960px); }}
    figure {{ margin: 0; background: rgba(255,255,255,.06); border: 1px solid rgba(255,255,255,.1); border-radius: 8px; overflow: hidden; }}
    img {{ display: block; width: 100%; aspect-ratio: 16 / 9; object-fit: contain; background: #020617; }}
    figcaption {{ padding: 10px 12px; color: #cbd5e1; font-size: 13px; }}
    details {{ margin-top: 16px; }}
    pre {{ overflow: auto; background: #020617; padding: 14px; border-radius: 8px; color: #d1d5db; }}
  </style>
</head>
<body>
  <header>
    <h1>PPT Image-Native Visual PoC</h1>
    <div class="meta">Stage: {stage} · PPTX: <a href="image_native_visual_poc.pptx">image_native_visual_poc.pptx</a> · Manifest: <a href="manifest.json">manifest.json</a> · Log: <a href="api_debug.log">api_debug.log</a></div>
  </header>
  {error_html}
  {''.join(rows)}
</body>
</html>"""
    (run_dir / "comparison.html").write_text(html, encoding="utf-8")


def run_edit_only(args: argparse.Namespace, env: Dict[str, str], run_dir: Path, log: RunLogger) -> int:
    image_path = Path(args.edit_image).expanduser() if args.edit_image else latest_existing_slide_artifact("*/cover_A.png")
    strategy = str(getattr(args, "edit_strategy", "masked")).strip().lower()
    mask_path = Path(args.edit_mask).expanduser() if args.edit_mask else image_path.with_name(image_path.name.replace("_A.png", "_mask.png"))
    if not image_path.exists():
        raise FileNotFoundError(f"edit image not found: {image_path}")
    if strategy == "masked" and not mask_path.exists():
        raise FileNotFoundError(f"edit mask not found: {mask_path}")
    if strategy != "masked" and not mask_path.exists():
        mask_path = None
    out_path = Path(args.edit_output).expanduser() if args.edit_output else run_dir / image_path.name.replace("_A.png", "_B_edit_only.png")
    if not out_path.is_absolute():
        out_path = run_dir / out_path
    log.log("edit_only_start", image=str(image_path), mask=str(mask_path) if mask_path is not None else "", output=str(out_path), edit_strategy=strategy)
    manifest: Dict[str, Any] = {
        "mode_requested": args.mode,
        "mode_effective": "api",
        "stage": "edit-only",
        "edit_strategy": strategy,
        "image": str(image_path),
        "mask": str(mask_path) if mask_path is not None else "",
        "output": str(out_path),
        "errors": [],
        "env_summary": env_summary(env),
    }
    try:
        manifest["edit"] = edit_image_api(image_path, mask_path, image_edit_prompt(strategy), out_path, env, log)
        normalize_slide_image(out_path)
        log.log("edit_only_done", output=str(out_path), bytes=out_path.stat().st_size)
    except Exception as exc:
        log.log("edit_only_failed", exception_type=type(exc).__name__, error=str(exc))
        manifest["errors"].append({"stage": "edit-only", "error": str(exc)})
    manifest["log"] = str(run_dir / "api_debug.log")
    manifest["api_events"] = log.events
    (run_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"run_dir": str(run_dir), "output": str(out_path) if out_path.exists() else "", "log": str(run_dir / "api_debug.log"), "errors": manifest["errors"]}, ensure_ascii=False, indent=2))
    return 2 if manifest["errors"] else 0


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--mode", choices=["auto", "api", "dry-run"], default="auto")
    parser.add_argument("--stage", choices=["full", "generate-only", "edit-only"], default="full")
    parser.add_argument("--slide-ids", default="", help="Comma-separated slide ids to run (e.g. metrics,process). Used by full/generate-only.")
    parser.add_argument("--max-slides", type=int, default=0)
    parser.add_argument("--vlm-provider", choices=["auto", "qwen", "gpt54"], default="auto")
    parser.add_argument("--vlm-model", default="")
    parser.add_argument("--vlm-endpoint", default="")
    parser.add_argument("--vlm-api-key", default="")
    parser.add_argument("--edit-strategy", choices=["masked", "prompt-only"], default="masked", help="B generation strategy: use mask or pure prompt.")
    parser.add_argument("--edit-image", default="", help="Existing A image for --stage edit-only. Defaults to latest runs/*/cover_A.png.")
    parser.add_argument("--edit-mask", default="", help="Existing visible text mask for --stage edit-only. Defaults to sibling *_mask.png.")
    parser.add_argument("--edit-output", default="", help="Output path for --stage edit-only. Defaults to current run dir.")
    args = parser.parse_args()

    env = load_env(ENV_PATH)
    schema = read_schema()
    run_dir = ensure_run_dir()
    log = RunLogger(run_dir / "api_debug.log")
    manifest: Dict[str, Any] = {
        "mode_requested": args.mode,
        "slides": {},
        "layouts": {},
        "errors": [],
        "api_events": [],
    }
    use_api = args.mode == "api" or (args.mode == "auto" and bool(env.get("AZURE_IMAGE_OPENAI_API_KEY", "").strip()))
    manifest["mode_effective"] = "api" if use_api else "dry-run"
    manifest["env_summary"] = env_summary(env)
    manifest["stage"] = args.stage
    manifest["slide_ids_requested"] = parse_csv_tokens(args.slide_ids)
    manifest["max_slides"] = args.max_slides
    resolved_vlm_provider = args.vlm_provider
    if resolved_vlm_provider == "auto":
        resolved_vlm_provider = "qwen"
    manifest["vlm_provider"] = resolved_vlm_provider
    manifest["vlm_model_override"] = args.vlm_model or ""
    manifest["vlm_endpoint_override"] = args.vlm_endpoint or ""
    manifest["edit_strategy"] = args.edit_strategy
    log.log(
        "run_start",
        run_dir=str(run_dir),
        mode_requested=args.mode,
        mode_effective=manifest["mode_effective"],
        stage=args.stage,
        slide_ids=manifest["slide_ids_requested"],
        edit_strategy=args.edit_strategy,
        max_slides=args.max_slides,
        vlm_provider=resolved_vlm_provider,
        vlm_model_override=args.vlm_model or "",
        vlm_endpoint_override=args.vlm_endpoint or "",
        env=manifest["env_summary"],
    )
    if args.stage == "edit-only":
        if not use_api:
            raise RuntimeError("--stage edit-only requires --mode api or AZURE_IMAGE_OPENAI_API_KEY")
        return run_edit_only(args, env, run_dir, log)
    background_paths: Dict[str, Path] = {}

    slides = list(schema.get("slides", []))
    requested_slide_ids = parse_csv_tokens(args.slide_ids)
    if requested_slide_ids:
        slide_by_id = {str(s.get("id")): s for s in slides}
        unknown = [sid for sid in requested_slide_ids if sid not in slide_by_id]
        if unknown:
            raise ValueError(f"--slide-ids contains unknown ids: {unknown}")
        slides = [slide_by_id[sid] for sid in requested_slide_ids]
    manifest["slide_ids_effective"] = [str(s.get("id")) for s in slides]
    if args.max_slides and args.max_slides > 0:
        slides = slides[: args.max_slides]

    for slide in slides:
        sid = safe_name(slide["id"])
        log.log("slide_start", slide=sid, type=slide.get("type", ""), text_slots=len(slide.get("text_slots", [])), chart_slots=len(slide.get("chart_slots", [])))
        a_path = run_dir / f"{sid}_A.png"
        mask_path = run_dir / f"{sid}_mask.png"
        b_path = run_dir / f"{sid}_B.png"
        c_preview_path = run_dir / f"{sid}_C_preview.png"
        try:
            if use_api:
                manifest["slides"][sid] = {"generation": generate_image_api(slide_prompt(slide), a_path, env, log)}
            else:
                log.log("image_generation_dry_run", slide=sid, out_path=str(a_path))
                dry_run_image(slide, a_path)
                manifest["slides"][sid] = {"generation": {"ok": True, "dry_run": True}}
        except Exception as exc:
            log.log("image_generation_fallback", slide=sid, exception_type=type(exc).__name__, error=str(exc), fallback="dry_run")
            manifest["errors"].append({"slide": sid, "stage": "generation", "error": str(exc)})
            dry_run_image(slide, a_path)
            manifest["slides"][sid] = {"generation": {"ok": False, "fallback": "dry_run", "error": str(exc)}}
        normalize_slide_image(a_path)
        log.log("image_A_ready", slide=sid, path=str(a_path), bytes=a_path.stat().st_size)

        if args.stage == "generate-only":
            continue

        layout = (
            infer_layout_with_vlm(
                a_path,
                slide,
                env,
                log,
                provider=resolved_vlm_provider,
                model_override=args.vlm_model,
                endpoint_override=args.vlm_endpoint,
                api_key_override=args.vlm_api_key,
            )
            if use_api
            else default_layout(slide)
        )
        manifest["layouts"][sid] = layout
        (run_dir / f"{sid}_layout.json").write_text(json.dumps(layout, ensure_ascii=False, indent=2), encoding="utf-8")
        build_mask(layout, mask_path)
        log.log("mask_ready", slide=sid, path=str(mask_path), bytes=mask_path.stat().st_size, layout_source=layout.get("source", ""))

        edit_prompt = image_edit_prompt(args.edit_strategy)
        try:
            if use_api:
                active_mask = mask_path if args.edit_strategy == "masked" else None
                manifest["slides"][sid]["edit"] = edit_image_api(a_path, active_mask, edit_prompt, b_path, env, log)
            else:
                log.log("image_edit_dry_run", slide=sid, out_path=str(b_path))
                dry_remove_text(a_path, mask_path, b_path)
                manifest["slides"][sid]["edit"] = {"ok": True, "dry_run": True}
        except Exception as exc:
            log.log("image_edit_fallback", slide=sid, exception_type=type(exc).__name__, error=str(exc), fallback="local_inpaint")
            manifest["errors"].append({"slide": sid, "stage": "edit", "error": str(exc)})
            local_inpaint_text(a_path, mask_path, b_path, log)
            manifest["slides"][sid]["edit"] = {"ok": False, "fallback": "local_inpaint", "error": str(exc)}
        normalize_slide_image(b_path)
        log.log("image_B_ready", slide=sid, path=str(b_path), bytes=b_path.stat().st_size)

        render_overlay_preview(b_path, slide, layout, c_preview_path)
        log.log("preview_ready", slide=sid, path=str(c_preview_path), bytes=c_preview_path.stat().st_size)
        background_paths[slide["id"]] = b_path

    pptx_path = run_dir / "image_native_visual_poc.pptx"
    if args.stage == "full":
        make_pptx(schema, manifest["layouts"], background_paths, pptx_path)
        log.log("pptx_ready", path=str(pptx_path), bytes=pptx_path.stat().st_size)
        manifest["pptx"] = str(pptx_path)
    else:
        manifest["pptx"] = ""
    manifest["html"] = str(run_dir / "comparison.html")
    manifest["log"] = str(run_dir / "api_debug.log")
    manifest["api_events"] = log.events
    (run_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    build_html(schema, run_dir, manifest)
    log.log("run_done", html=str(run_dir / "comparison.html"), pptx=str(pptx_path) if args.stage == "full" else "", errors=len(manifest["errors"]))
    manifest["api_events"] = log.events
    (run_dir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({"run_dir": str(run_dir), "html": str(run_dir / "comparison.html"), "pptx": str(pptx_path) if args.stage == "full" else "", "log": str(run_dir / "api_debug.log")}, ensure_ascii=False, indent=2))
    if args.mode == "api" and manifest["errors"]:
        print("API mode completed with fallback artifacts; see manifest.json for API errors.", file=sys.stderr)
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
